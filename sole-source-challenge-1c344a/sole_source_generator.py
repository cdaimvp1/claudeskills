"""
sole_source_generator.py
Lilly Procurement Skills, sole-source-challenge-1c344a deterministic
PPTX/DOCX generator.

Purpose (mirrors should-cost-builder-1c344a/should_cost_generator.py and
evaluation-engine-1c344a/evaluation_report_generator.py's "generator
scripts, not freehand authoring" pattern): this module takes a validated
sole-source challenge register (the request, the seven scored challenge
dimensions, mitigations, the alternatives register, price-validation
evidence, the research log, and SME routing, shaped per
sole-source-challenge-1c344a/SKILL.md's Phase 2-5 and Deliverables
sections) and MECHANICALLY produces the sole_source_challenge_report in
EITHER of two output formats from the SAME data model:
  - a PowerPoint deck (.pptx, python-pptx): one idea per slide, title plus
    concise bullets/tables.
  - a Word memo (.docx, python-docx): the full narrative memo, Sections
    01-06 plus an Appendix, exactly the skeleton SKILL.md's Deliverable #1
    describes ("Same skeleton every run regardless of mode or verdict;
    only the content and the 04 branch content vary.").

Two-stage discipline (same as should-cost/evaluation-engine):
  1. The Sole-Source Defensibility Score is computed ONLY by calling the
     vendored numeric_kernel.py's weighted_score(), never by model or
     script arithmetic. Per SKILL.md's G11 HARD RULE ("This skill vendors
     numeric_kernel.py. The Sole-Source Defensibility Score is computed
     ONLY by calling weighted_score() in that kernel, never by model
     arithmetic.") and Phase 4's "Computation requirement (HARD RULE): do
     not hand-compute the total."
  2. The seven dimension weights (0.20/0.20/0.15/0.15/0.10/0.10/0.10) are
     FIXED by SKILL.md's Phase 4 table, not caller-supplied. This module
     hardcodes them in DIMENSION_SPECS below (JUDGMENT CALL, flagged: an
     input register that includes a "weight" key is validated against
     that fixed constant and REJECTED on mismatch, rather than silently
     trusting a caller-supplied number, since Phase 4 states the weights
     "sum to 1.00" as a closed, fixed table, unlike evaluation-engine's
     genuinely per-run configurable criterion weights).

Hard-coded invariants (code-level checks, not comments; mirrors
should-cost's run_hardcoded_invariant_checks() / evaluation-engine's
identically-named function):
  - Weights-sum-to-one invariant (defense in depth on the hardcoded
    constant itself; Phase 4: "0.20+0.20+0.15+0.15+0.10+0.10+0.10 = 1.00").
  - Weighted-contributions-reconcile invariant: every dimension's
    score x weight must sum exactly to the kernel-computed
    defensibility_score (mirrors evaluation-engine's Canonical Dashboard
    "Numbers-reconcile assertion").
  - Defensibility-score-range invariant: 0.0-5.0.
  - Verdict-matches-score-band-and-cap invariant: the verdict band
    (Phase 5 table) must match the score, AND Hard Rule 2's "if the
    majority of scored dimensions are ASSERTED rather than VERIFIED, cap
    the verdict at DEFENSIBLE WITH MITIGATIONS even if the raw score would
    read DEFENSIBLE" must be applied exactly when it applies, never more,
    never less.
  - Price-position-consistent invariant: Rule 5 ("Price-validation cannot
    be scored on faith") - when a method other than "none" is supplied,
    the stated within/above/below position must match the actual
    band_low/band_high/sole_source_price numbers.
  - Mitigations-required-when-conditioned invariant: a DEFENSIBLE WITH
    MITIGATIONS verdict must carry at least one mitigation with a named
    owner and due point (Phase 5: "Each mitigation names an owner and a
    due point").
  - Alternatives-register-nonempty invariant: Deliverable #4's "If no
    alternatives were evaluated ... emit the file with a single row
    stating 'no alternatives evaluated' and why" - alternatives is never
    truly empty.
  - Ranked-alternatives-sorted invariant: the WEAK-verdict "ranked
    alternatives_register.csv entries worth pursuing" table is genuinely
    ranked (viability tier, then confidence, then name), not just listed
    in input order.
  - Gating-items-routed invariant: Hard Rule 6 ("Gating items are
    flagged, never cleared") - if the validated input text references a
    debarment/sanctions/GxP-adjacent keyword, sme_routing must be
    non-empty.
  - No-em-dash invariant (validation time): the suite-wide HARD RULE
    ("Never use em dashes... in ANY written output") is enforced by
    scanning every string leaf of the raw register before it is trusted.

If any invariant fails, this script RAISES rather than writing a document
that fails its own reconciliation.

Scope note: this script builds the sole_source_challenge_report generator
(both formats) ONLY, per SKILL.md Deliverable #1. The other mandated
artifacts (challenge_scorecard.csv, alternatives_register.csv,
sole_source_justification_handoff.json, and the four-tab canonical
dashboard) are separate deliverables and are NOT produced here.

============================================================================
Sole-source register input contract (JUDGMENT CALL, flagged): SKILL.md
gives the seven-dimension table, the verdict bands, and the handoff JSON
shape, but not a single raw INPUT contract for a report generator. This
module resolves that ambiguity as follows.

Narrative provenance (JUDGMENT CALL, flagged): mitigations, the
alternatives register, price-validation figures, the research log, and
SME routing are this skill's own upstream analytical judgment calls
(Phases 2-4), not something a deterministic document renderer should
invent. This generator therefore accepts them as already-decided input,
the same "consume, don't re-derive" discipline evaluation-engine applies
to its recommendation.key_tradeoffs / secondary_tradeoffs fields. Only the
defensibility_score, the verdict (including the Hard Rule 2 cap), the
weighted contributions, the ranked alternatives ordering, the
recommended_next_action, and the closing Next Steps are COMPUTED by this
generator, never accepted as input, mirroring evaluation-engine's
"recommended supplier ... COMPUTED by this generator, never accepted as
input" pattern.

recommended_next_action and next_steps are DERIVED, not accepted as raw
input (JUDGMENT CALL, flagged): Phase 5's "Recommended next action (WEAK
verdict), by mode" table and the "Next Steps (closing template)" section
both give fixed, quotable template text keyed off verdict/mode, so this
generator reproduces that text mechanically from already-validated
ground truth rather than accepting free-text from the caller (which would
risk a stale or inconsistent narrative sneaking past the kernel-computed
verdict).

Research-pending threshold (JUDGMENT CALL, flagged): G7 states a 2-search
minimum for the Phase 3 "light market-check" specifically, not for all
research broadly. This generator does not distinguish research-log entry
types, so it applies the len(research_log) < 2 test to the WHOLE research
log as a conservative proxy, and labels the result "RESEARCH PENDING"
accordingly. This is flagged here and in the rendered output.

Register shape:
{
  "request": {
    "supplier": "string (MUST/blocking)",
    "need_description": "string (MUST/blocking)",
    "stated_rationale": "string (MUST/blocking, the requester's own stated reason)",
    "mode": "NEW" | "RENEWAL" | "AUDIT",
    "date": "string (as-of date)",
    "category": "string (optional)",
    "requester": "string (optional)",
    "est_value_usd": float | null (optional),
    "term": "string (optional)"
  },
  "dimensions": [
    {"dimension_key": "unique_capability", "score": 0.0-5.0, "label": "VERIFIED"|"ASSERTED"|"INFERRED",
     "confidence": "HIGH"|"MEDIUM"|"LOW", "evidence": "string", "rationale": "string"},
    ...  # exactly 7, one per fixed key: unique_capability, constraint_basis, competition_history,
    #      requirements_separability, alt_availability, urgency_legitimacy, price_validation
  ],
  "mitigations": [{"dimension": "string", "action": "string", "owner": "string", "due": "string"}, ...],
  "alternatives": [
    {"candidate_name": "string", "origin": "supplier-landscape-excluded"|"market-check"|"user-provided"|"aria-shared-incumbent-history",
     "original_exclusion_reason": "string (required if origin==supplier-landscape-excluded)",
     "capability_gap": "string", "reassessed_viability": "Not Viable"|"Viable With Gaps"|"Viable",
     "confidence": "HIGH"|"MEDIUM"|"LOW", "source": "string", "date": "string"}, ...
  ],  # non-empty; a single "no alternatives evaluated" row is valid
  "price_validation": {"method": "should-cost"|"market-rate"|"prior-rate"|"rate-card"|"none",
                        "source": "string", "band_low": float|null, "band_high": float|null,
                        "sole_source_price": float|null, "position": "within"|"above"|"below"|"not_computable"},
  "research_log": [{"query": "string", "source": "string", "date": "string", "results": int}, ...],
  "consumed_artifacts": ["string", ...]  (optional),
  "sme_routing": [{"issue": "string", "route_to": "string", "reason": "string"}, ...],
  "research_methodology_note": "string (optional)"
}
============================================================================
"""

from __future__ import annotations

import copy
import os
import sys
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Sequence, Tuple

# ---------------------------------------------------------------------------
# Vendored numeric kernel (same directory). Per SKILL.md's G11 HARD RULE and
# Phase 4's "Computation requirement": the Defensibility Score MUST be
# computed by calling weighted_score() in the vendored numeric_kernel.py,
# never hand-computed.
# ---------------------------------------------------------------------------
_KERNEL_IMPORT_ERROR: Optional[Exception] = None
try:
    _THIS_DIR = os.path.dirname(os.path.abspath(__file__))
    if _THIS_DIR not in sys.path:
        sys.path.insert(0, _THIS_DIR)
    from numeric_kernel import weighted_score as kernel_weighted_score
    from numeric_kernel import WeightSumError as KernelWeightSumError
    from numeric_kernel import InvalidInputError as KernelInvalidInputError
    KERNEL_AVAILABLE = True
except Exception as _exc:  # pragma: no cover - defensive, disclosed at runtime
    KERNEL_AVAILABLE = False
    _KERNEL_IMPORT_ERROR = _exc

    def kernel_weighted_score(*args, **kwargs):  # type: ignore
        raise RuntimeError(
            "numeric_kernel.py unavailable; cannot compute the kernel-required "
            f"Defensibility Score. Import error: {_KERNEL_IMPORT_ERROR}"
        )

    class KernelWeightSumError(Exception):  # type: ignore
        pass

    class KernelInvalidInputError(Exception):  # type: ignore
        pass

# ---------------------------------------------------------------------------
# python-docx detection. Mirrors evaluation_report_generator.py: try the
# import; if unavailable, raise a clear ImportError at document-build time
# (not at module-import time), so validation and ground-truth logic (which
# need no DOCX/PPTX library) remain usable regardless.
# ---------------------------------------------------------------------------
try:
    import docx  # noqa: F401
    from docx import Document
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn as docx_qn
    from docx.oxml import OxmlElement
    DOCX_AVAILABLE = True
    DOCX_VERSION = getattr(docx, "__version__", "unknown")
except Exception as _exc:  # pragma: no cover
    DOCX_AVAILABLE = False
    DOCX_VERSION = None
    _DOCX_IMPORT_ERROR = _exc

# ---------------------------------------------------------------------------
# python-pptx detection.
# ---------------------------------------------------------------------------
try:
    import pptx  # noqa: F401
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
    PPTX_VERSION = getattr(pptx, "__version__", "unknown")
except Exception as _exc:  # pragma: no cover
    PPTX_AVAILABLE = False
    PPTX_VERSION = None
    _PPTX_IMPORT_ERROR = _exc


class SoleSourceValidationError(Exception):
    """Raised when the sole-source register is missing a required field or
    carries a value the skill's Hard Rules forbid. Per Hard Rule 7 ("No
    fabricated evidence, no boilerplate filler, no em dashes") this module
    refuses rather than guessing a missing or malformed value."""


class ReconciliationError(Exception):
    """Raised when a hard-coded invariant fails. The document is not saved
    when this fires."""


# ===========================================================================
# House style constants (quoted from SKILL.md's "House style and palette"
# block: "Canonical non-green status palette only: positive text Bold Blue
# #0F3A85 on Neutral Sky #D4E5F7; warning text Amber #B45309 on Neutral
# Cream #FFF0D8; negative text Lilly Red #E1251B on Neutral Rose #FDE8E5;
# neutral/N-A Bold Grey #8A969E ... cards/borders Neutral Stone #E4EBF1;
# header bar Lilly Black #212121. No green or teal anywhere." Hardcoded
# here exactly as evaluation_report_generator.py hardcodes its own palette
# from lilly-brand-assets, keeping this generator self-contained.)
# ===========================================================================
LILLY_RED = "E1251B"
BOLD_BLUE = "0F3A85"
LILLY_BLACK = "212121"
AMBER = "B45309"
MUTED_GREY = "8A969E"
WHITE = "FFFFFF"
NEUTRAL_SKY = "D4E5F7"
NEUTRAL_CREAM = "FFF0D8"
NEUTRAL_ROSE = "FDE8E5"
NEUTRAL_STONE = "E4EBF1"

VALID_MODE = ("NEW", "RENEWAL", "AUDIT")
VALID_LABEL = ("VERIFIED", "ASSERTED", "INFERRED")
VALID_CONFIDENCE = ("HIGH", "MEDIUM", "LOW")
VALID_ORIGIN = ("supplier-landscape-excluded", "market-check", "user-provided", "aria-shared-incumbent-history")
VALID_VIABILITY = ("Not Viable", "Viable With Gaps", "Viable")
VALID_PRICE_METHOD = ("should-cost", "market-rate", "prior-rate", "rate-card", "none")
VALID_PRICE_POSITION = ("within", "above", "below", "not_computable")

GATING_KEYWORDS = ("debarment", "sanction", "gxp", "denied party", "excluded party")

MODE_DEFINITION = {
    "NEW": "the supplier has not yet been awarded; this challenge runs before the sole-source case is submitted for approval.",
    "RENEWAL": "an existing sole-source arrangement is up for its periodic re-justification.",
    "AUDIT": "a third party (Compliance, procurement leadership, an internal audit) wants an independent challenge of an arrangement they did not originate.",
}


# ===========================================================================
# Fixed dimension table (Phase 4, "Dimensions and weights (fixed; sum to
# 1.00)"). Weights are NOT caller-supplied; see module docstring.
# ===========================================================================

@dataclass(frozen=True)
class DimensionSpec:
    key: str
    order: int
    label: str
    weight: float
    anchor_high: str
    anchor_low: str


DIMENSION_SPECS: List[DimensionSpec] = [
    DimensionSpec(
        "unique_capability", 1, "Unique Capability", 0.20,
        "A specific, evidenced capability/technology/certification no identified alternative holds, tied to a stated requirement",
        "The capability is common; multiple suppliers plausibly offer it",
    ),
    DimensionSpec(
        "constraint_basis", 2, "Constraint Basis", 0.20,
        "A named, evidenced constraint (existing IP ownership, a safety-validated system, a regulatory qualification, a continuity dependency) structurally locks in this supplier",
        "No named constraint; supplier preference asserted without a structural reason",
    ),
    DimensionSpec(
        "competition_history", 3, "Competition History", 0.15,
        "A real competitive process was run and this supplier won, OR a documented, dated market scan found no viable alternative",
        "No competitive process was run or attempted, and none is evidenced",
    ),
    DimensionSpec(
        "requirements_separability", 4, "Requirements Separability", 0.15,
        "Requirements were reviewed and genuinely cannot be unbundled without breaking the constraint",
        "Requirements could plausibly be split to open at least part of the scope, but were not, or were not even reviewed",
    ),
    DimensionSpec(
        "alt_availability", 5, "Alternative Supplier Availability", 0.10,
        "A real check (supplier-landscape output or this skill's own market-check) found no viable alternative",
        "Viable alternatives were identified and not seriously considered",
    ),
    DimensionSpec(
        "urgency_legitimacy", 6, "Urgency Legitimacy", 0.10,
        "Urgency traces to an external driver (regulatory deadline, safety event, a contract cliff not created by the business)",
        "Urgency is self-created (late planning, a deferred renewal, poor lead time)",
    ),
    DimensionSpec(
        "price_validation", 7, "Price Validation Substitute", 0.10,
        "A should-cost model or market-rate benchmark independently validates the price despite no competitive bid",
        "No price validation attempted; the price is taken on faith",
    ),
]
DIMENSION_KEYS: List[str] = [d.key for d in DIMENSION_SPECS]
DIMENSION_BY_KEY: Dict[str, DimensionSpec] = {d.key: d for d in DIMENSION_SPECS}
FIXED_WEIGHTS: Dict[str, float] = {d.key: d.weight for d in DIMENSION_SPECS}

VERDICT_LABEL = {
    "DEFENSIBLE": "Defensible",
    "DEFENSIBLE_WITH_MITIGATIONS": "Defensible With Mitigations",
    "WEAK_RECOMMEND_COMPETITION": "Weak, Recommend Competitive Alternative",
}
VERDICT_COLOR = {
    "DEFENSIBLE": (BOLD_BLUE, NEUTRAL_SKY),
    "DEFENSIBLE_WITH_MITIGATIONS": (AMBER, NEUTRAL_CREAM),
    "WEAK_RECOMMEND_COMPETITION": (LILLY_RED, NEUTRAL_ROSE),
}
_VIAB_RANK = {"Viable": 3, "Viable With Gaps": 2, "Not Viable": 1}
_CONF_RANK = {"HIGH": 3, "MEDIUM": 2, "LOW": 1}


# ===========================================================================
# 1. Sole-source register: typed input
# ===========================================================================

@dataclass
class RequestInfo:
    supplier: str
    need_description: str
    stated_rationale: str
    mode: str
    date: str
    category: str
    requester: str
    est_value_usd: Optional[float]
    term: str


@dataclass
class DimensionEntry:
    key: str
    label_text: str
    weight: float
    score: float
    evidence_label: str
    confidence: str
    evidence: str
    rationale: str


@dataclass
class Mitigation:
    dimension: str
    action: str
    owner: str
    due: str
    auto_generated: bool = False


@dataclass
class AlternativeCandidate:
    candidate_name: str
    origin: str
    original_exclusion_reason: str
    capability_gap: str
    reassessed_viability: str
    confidence: str
    source: str
    date: str


@dataclass
class PriceValidationInfo:
    method: str
    source: str
    band_low: Optional[float]
    band_high: Optional[float]
    sole_source_price: Optional[float]
    position: str


@dataclass
class ResearchLogEntry:
    query: str
    source: str
    date: str
    results: int


@dataclass
class SmeRoutingEntry:
    issue: str
    route_to: str
    reason: str


@dataclass
class SoleSourceInput:
    request: RequestInfo
    dimensions: List[DimensionEntry]
    mitigations: List[Mitigation]
    alternatives: List[AlternativeCandidate]
    price_validation: PriceValidationInfo
    research_log: List[ResearchLogEntry]
    consumed_artifacts: List[str]
    sme_routing: List[SmeRoutingEntry]
    research_methodology_note: str


def _find_em_dashes(obj: Any, path: str = "root") -> List[str]:
    """Recursively scan every string leaf for the em dash character
    (U+2014), per the suite-wide HARD RULE ('Never use em dashes... in ANY
    written output')."""
    found: List[str] = []
    em_dash = chr(0x2014)
    if isinstance(obj, str):
        if em_dash in obj:
            found.append(path)
    elif isinstance(obj, dict):
        for k, v in obj.items():
            found.extend(_find_em_dashes(v, f"{path}.{k}"))
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            found.extend(_find_em_dashes(v, f"{path}[{i}]"))
    return found


def validate_sole_source_input(register: Dict[str, Any]) -> SoleSourceInput:
    """Validate a raw sole-source register dict and return a typed
    SoleSourceInput. Refuses (raises SoleSourceValidationError) rather than
    guessing when a required field is missing or a value violates a
    documented, fixed convention, per Hard Rule 7 ('No fabricated evidence,
    no boilerplate filler, no em dashes')."""
    errors: List[str] = []

    em_dash_paths = _find_em_dashes(register)
    if em_dash_paths:
        errors.append(
            "Em dash character (U+2014) found in the following field(s), violating the suite-wide "
            f"HARD RULE ('Never use em dashes'): {em_dash_paths}. Restructure with hyphens, colons, "
            "parentheses, or separate sentences instead."
        )

    # --- Request -------------------------------------------------------
    raw_request = register.get("request")
    if not isinstance(raw_request, dict):
        errors.append(
            "'request' must be an object with supplier, need_description, stated_rationale, mode, "
            "date (MUST/blocking inputs per SUITE INTERACTION PROTOCOL S5)."
        )
        raw_request = {}
    required_request = ["supplier", "need_description", "stated_rationale", "mode", "date"]
    missing_request = [f for f in required_request if f not in raw_request or not str(raw_request.get(f, "")).strip()]
    if missing_request:
        errors.append(f"'request' is missing or has empty required (blocking, per S5) field(s): {missing_request}.")

    supplier = str(raw_request.get("supplier", "")).strip()
    need_description = str(raw_request.get("need_description", "")).strip()
    stated_rationale = str(raw_request.get("stated_rationale", "")).strip()
    mode = str(raw_request.get("mode", "")).strip().upper()
    if mode not in VALID_MODE:
        errors.append(f"'request.mode' must be one of {VALID_MODE}; got {raw_request.get('mode')!r}.")
        mode = "NEW"
    date = str(raw_request.get("date", "")).strip()
    category = str(raw_request.get("category", "")).strip() or "NEEDS_INPUT, no category provided"
    requester = str(raw_request.get("requester", "")).strip() or "NEEDS_INPUT, no requester named"
    term = str(raw_request.get("term", "")).strip() or "NEEDS_INPUT, no term provided"
    est_value_raw = raw_request.get("est_value_usd")
    if est_value_raw is None:
        est_value_usd = None
    elif isinstance(est_value_raw, (int, float)) and not isinstance(est_value_raw, bool):
        est_value_usd = float(est_value_raw)
    else:
        errors.append(f"'request.est_value_usd', when present, must be numeric or null; got {est_value_raw!r}.")
        est_value_usd = None

    # --- Dimensions (exactly 7, fixed keys, fixed weights) --------------
    raw_dimensions = register.get("dimensions", [])
    if not isinstance(raw_dimensions, list) or len(raw_dimensions) != 7:
        got = len(raw_dimensions) if isinstance(raw_dimensions, list) else type(raw_dimensions).__name__
        errors.append(
            f"'dimensions' must be a list of exactly 7 entries, one per fixed dimension_key "
            f"{DIMENSION_KEYS} (SKILL.md Phase 4); got {got}."
        )
        raw_dimensions = []

    dims_by_key: Dict[str, DimensionEntry] = {}
    seen_keys: List[str] = []
    for i, d in enumerate(raw_dimensions):
        if not isinstance(d, dict):
            errors.append(f"'dimensions[{i}]' must be an object.")
            continue
        required = ["dimension_key", "score", "label", "confidence"]
        missing = [f for f in required if f not in d]
        if missing:
            errors.append(f"'dimensions[{i}]' is missing required field(s): {missing}.")
            continue
        key = str(d["dimension_key"]).strip()
        if key not in DIMENSION_BY_KEY:
            errors.append(f"'dimensions[{i}]' dimension_key {key!r} is not one of the 7 fixed dimensions {DIMENSION_KEYS}.")
            continue
        if key in seen_keys:
            errors.append(f"Duplicate dimension_key {key!r} in 'dimensions'.")
            continue
        seen_keys.append(key)
        spec = DIMENSION_BY_KEY[key]

        if "weight" in d:
            try:
                w = float(d["weight"])
                if abs(w - spec.weight) > 1e-9:
                    errors.append(
                        f"'dimensions[{i}]' ({key}) supplies weight={w}, which disagrees with the fixed "
                        f"weight {spec.weight} from SKILL.md Phase 4. Dimension weights are FIXED, not "
                        "caller-supplied; omit 'weight' or match it exactly."
                    )
            except (TypeError, ValueError):
                errors.append(f"'dimensions[{i}]' ({key}) 'weight', when present, must be numeric.")

        score = d["score"]
        if not isinstance(score, (int, float)) or isinstance(score, bool) or not (0.0 <= float(score) <= 5.0):
            errors.append(f"'dimensions[{i}]' ({key}) score must be numeric in [0.0, 5.0]; got {score!r}.")
            score = 0.0
        label = str(d["label"]).strip().upper()
        if label not in VALID_LABEL:
            errors.append(f"'dimensions[{i}]' ({key}) label must be one of {VALID_LABEL}; got {d['label']!r}.")
            label = "ASSERTED"
        confidence = str(d["confidence"]).strip().upper()
        if confidence not in VALID_CONFIDENCE:
            errors.append(f"'dimensions[{i}]' ({key}) confidence must be one of {VALID_CONFIDENCE}; got {d['confidence']!r}.")
            confidence = "LOW"
        evidence = str(d.get("evidence", "")).strip()
        if not evidence:
            evidence = "Not available (gap, per Hard Rule 7)."
        rationale = str(d.get("rationale", "")).strip()
        if not rationale:
            errors.append(
                f"'dimensions[{i}]' ({key}) 'rationale' must be a non-empty string tying the score to "
                "evidence (Hard Rule 7)."
            )
        dims_by_key[key] = DimensionEntry(
            key=key, label_text=spec.label, weight=spec.weight, score=float(score),
            evidence_label=label, confidence=confidence, evidence=evidence, rationale=rationale,
        )

    missing_keys = [k for k in DIMENSION_KEYS if k not in dims_by_key]
    if not errors and missing_keys:
        errors.append(
            f"'dimensions' is missing entries for: {missing_keys}. All 7 fixed dimensions are required "
            "every run (SKILL.md Hard Rule 8, 'Deliverable structure is deterministic')."
        )
    dimensions = [dims_by_key[k] for k in DIMENSION_KEYS if k in dims_by_key]

    # --- Mitigations -----------------------------------------------------
    raw_mitigations = register.get("mitigations", [])
    if not isinstance(raw_mitigations, list):
        errors.append("'mitigations' must be a list (may be empty).")
        raw_mitigations = []
    mitigations: List[Mitigation] = []
    for i, m in enumerate(raw_mitigations):
        if not isinstance(m, dict):
            errors.append(f"'mitigations[{i}]' must be an object.")
            continue
        required = ["dimension", "action", "owner", "due"]
        missing = [f for f in required if f not in m or not str(m.get(f, "")).strip()]
        if missing:
            errors.append(f"'mitigations[{i}]' is missing or has empty required field(s): {missing}.")
            continue
        mitigations.append(Mitigation(
            dimension=str(m["dimension"]).strip(), action=str(m["action"]).strip(),
            owner=str(m["owner"]).strip(), due=str(m["due"]).strip(),
        ))

    # --- Alternatives ------------------------------------------------------
    raw_alternatives = register.get("alternatives", [])
    if not isinstance(raw_alternatives, list) or len(raw_alternatives) == 0:
        errors.append(
            "'alternatives' must be a non-empty list; if no candidate was genuinely evaluated, supply a "
            "single row explaining why (Deliverable #4, 'no alternatives evaluated')."
        )
        raw_alternatives = []
    alternatives: List[AlternativeCandidate] = []
    for i, a in enumerate(raw_alternatives):
        if not isinstance(a, dict):
            errors.append(f"'alternatives[{i}]' must be an object.")
            continue
        required = ["candidate_name", "origin", "capability_gap", "reassessed_viability", "confidence", "source", "date"]
        missing = [f for f in required if f not in a]
        if missing:
            errors.append(f"'alternatives[{i}]' is missing required field(s): {missing}.")
            continue
        origin = str(a["origin"]).strip()
        if origin not in VALID_ORIGIN:
            errors.append(f"'alternatives[{i}]' origin must be one of {VALID_ORIGIN}; got {origin!r}.")
            continue
        viability = str(a["reassessed_viability"]).strip()
        if viability not in VALID_VIABILITY:
            errors.append(f"'alternatives[{i}]' reassessed_viability must be one of {VALID_VIABILITY}; got {viability!r}.")
            continue
        confidence = str(a["confidence"]).strip().upper()
        if confidence not in VALID_CONFIDENCE:
            errors.append(f"'alternatives[{i}]' confidence must be one of {VALID_CONFIDENCE}; got {a['confidence']!r}.")
            continue
        excl_reason = str(a.get("original_exclusion_reason", "")).strip()
        if origin == "supplier-landscape-excluded" and not excl_reason:
            errors.append(
                f"'alternatives[{i}]' origin is 'supplier-landscape-excluded' but "
                "'original_exclusion_reason' is empty; the original exclusion reason must carry "
                "forward (SKILL.md Phase 3)."
            )
        if not excl_reason:
            excl_reason = "n/a, newly surfaced this run"
        alternatives.append(AlternativeCandidate(
            candidate_name=str(a["candidate_name"]).strip(), origin=origin,
            original_exclusion_reason=excl_reason, capability_gap=str(a["capability_gap"]).strip(),
            reassessed_viability=viability, confidence=confidence,
            source=str(a["source"]).strip(), date=str(a["date"]).strip(),
        ))

    # --- Price validation --------------------------------------------------
    raw_pv = register.get("price_validation")
    if not isinstance(raw_pv, dict):
        errors.append(
            "'price_validation' must be an object (Rule 5, 'Price-validation cannot be scored on faith')."
        )
        raw_pv = {}
    method = str(raw_pv.get("method", "")).strip().lower()
    if method not in VALID_PRICE_METHOD:
        errors.append(f"'price_validation.method' must be one of {VALID_PRICE_METHOD}; got {raw_pv.get('method')!r}.")
        method = "none"
    position = str(raw_pv.get("position", "")).strip().lower()
    if position not in VALID_PRICE_POSITION:
        errors.append(f"'price_validation.position' must be one of {VALID_PRICE_POSITION}; got {raw_pv.get('position')!r}.")
        position = "not_computable"
    band_low = raw_pv.get("band_low")
    band_high = raw_pv.get("band_high")
    sole_source_price = raw_pv.get("sole_source_price")
    source = str(raw_pv.get("source", "")).strip()

    if method == "none":
        if band_low is not None or band_high is not None or sole_source_price is not None:
            errors.append(
                "'price_validation.method' is 'none' but band_low/band_high/sole_source_price are "
                "populated; either supply a real method or clear the figures."
            )
        if position != "not_computable":
            errors.append("'price_validation.method' is 'none', so 'position' must be 'not_computable'.")
        band_low = band_high = sole_source_price = None
    else:
        if not source:
            errors.append(f"'price_validation.source' must be a non-empty citation when method is {method!r} (Rule 5).")
        for fname, fval in (("band_low", band_low), ("band_high", band_high), ("sole_source_price", sole_source_price)):
            if not isinstance(fval, (int, float)) or isinstance(fval, bool):
                errors.append(f"'price_validation.{fname}' must be numeric when method is {method!r}; got {fval!r}.")
        if (isinstance(band_low, (int, float)) and not isinstance(band_low, bool)
                and isinstance(band_high, (int, float)) and not isinstance(band_high, bool)
                and band_low > band_high):
            errors.append(f"'price_validation.band_low' ({band_low}) must be <= 'band_high' ({band_high}).")
        if position == "not_computable":
            errors.append(
                f"'price_validation.position' is 'not_computable' but method is {method!r} with figures "
                "supplied; state within/above/below."
            )

    price_validation = PriceValidationInfo(
        method=method, source=source,
        band_low=float(band_low) if isinstance(band_low, (int, float)) and not isinstance(band_low, bool) else None,
        band_high=float(band_high) if isinstance(band_high, (int, float)) and not isinstance(band_high, bool) else None,
        sole_source_price=(
            float(sole_source_price)
            if isinstance(sole_source_price, (int, float)) and not isinstance(sole_source_price, bool)
            else None
        ),
        position=position,
    )

    # --- Research log --------------------------------------------------
    raw_rl = register.get("research_log", [])
    if not isinstance(raw_rl, list):
        errors.append("'research_log' must be a list (may be empty).")
        raw_rl = []
    research_log: List[ResearchLogEntry] = []
    for i, r in enumerate(raw_rl):
        if not isinstance(r, dict):
            errors.append(f"'research_log[{i}]' must be an object.")
            continue
        required = ["query", "source", "date", "results"]
        missing = [f for f in required if f not in r]
        if missing:
            errors.append(f"'research_log[{i}]' is missing field(s): {missing}.")
            continue
        results = r["results"]
        if not isinstance(results, int) or isinstance(results, bool) or results < 0:
            errors.append(f"'research_log[{i}].results' must be a non-negative integer; got {results!r}.")
            results = 0
        research_log.append(ResearchLogEntry(
            query=str(r["query"]).strip(), source=str(r["source"]).strip(),
            date=str(r["date"]).strip(), results=int(results),
        ))

    consumed_artifacts = [str(x).strip() for x in register.get("consumed_artifacts", []) if str(x).strip()]

    # --- SME routing -------------------------------------------------------
    raw_sme = register.get("sme_routing", [])
    if not isinstance(raw_sme, list):
        errors.append("'sme_routing' must be a list (may be empty).")
        raw_sme = []
    sme_routing: List[SmeRoutingEntry] = []
    for i, s in enumerate(raw_sme):
        if not isinstance(s, dict):
            errors.append(f"'sme_routing[{i}]' must be an object.")
            continue
        required = ["issue", "route_to", "reason"]
        missing = [f for f in required if f not in s or not str(s.get(f, "")).strip()]
        if missing:
            errors.append(f"'sme_routing[{i}]' is missing or has empty required field(s): {missing}.")
            continue
        sme_routing.append(SmeRoutingEntry(
            issue=str(s["issue"]).strip(), route_to=str(s["route_to"]).strip(), reason=str(s["reason"]).strip(),
        ))

    research_methodology_note = str(register.get("research_methodology_note", "")).strip()
    if not research_methodology_note:
        n = len(research_log)
        research_methodology_note = (
            f"No research methodology note was supplied for this run. {n} "
            f"entr{'y' if n == 1 else 'ies'} logged; see the Research Log below and the G7 minimum "
            "compliance check in Section 06."
        )

    if errors:
        raise SoleSourceValidationError(
            "Sole-source register failed validation; refusing to guess missing or invalid fields. "
            "Issues found:\n  - " + "\n  - ".join(errors)
        )

    return SoleSourceInput(
        request=RequestInfo(
            supplier=supplier, need_description=need_description, stated_rationale=stated_rationale,
            mode=mode, date=date, category=category, requester=requester,
            est_value_usd=est_value_usd, term=term,
        ),
        dimensions=dimensions, mitigations=mitigations, alternatives=alternatives,
        price_validation=price_validation, research_log=research_log,
        consumed_artifacts=consumed_artifacts, sme_routing=sme_routing,
        research_methodology_note=research_methodology_note,
    )


# ===========================================================================
# Ground-truth computation (via the vendored kernel only, per G11 HARD RULE)
# ===========================================================================

def _rank_key(a: AlternativeCandidate) -> Tuple[int, int, str]:
    return (-_VIAB_RANK[a.reassessed_viability], -_CONF_RANK[a.confidence], a.candidate_name)


@dataclass
class GroundTruth:
    weights: Dict[str, float]
    weighted_contribution: Dict[str, float]
    defensibility_score: float
    raw_verdict: str
    label_counts: Dict[str, int]
    asserted_majority: bool
    verdict: str
    verdict_capped: bool
    verdict_label: str
    verdict_color: str
    verdict_bg: str
    weakest: DimensionEntry
    strongest: List[DimensionEntry]
    effective_mitigations: List[Mitigation]
    ranked_alternatives: List[AlternativeCandidate]
    research_pending: bool
    recommended_next_action: str
    next_steps: List[str]


def compute_ground_truth(reg: SoleSourceInput) -> GroundTruth:
    """Compute every derived figure this report shows, calling ONLY
    numeric_kernel.weighted_score() for the Defensibility Score, per
    SKILL.md's G11 HARD RULE. This is the reference against which the
    hard-coded invariants are checked before the document is saved."""
    if not KERNEL_AVAILABLE:
        raise RuntimeError(
            "numeric_kernel.py could not be imported; this generator refuses to hand-compute the "
            f"Defensibility Score in its place (SKILL.md G11 HARD RULE). Import error: {_KERNEL_IMPORT_ERROR}"
        )

    scores = {d.key: d.score for d in reg.dimensions}
    weights = dict(FIXED_WEIGHTS)
    try:
        score = kernel_weighted_score(scores, weights)
    except KernelWeightSumError as e:
        raise ReconciliationError(
            f"Weights-sum-to-one invariant FAILED: numeric_kernel.weighted_score() refused because the "
            f"fixed dimension weights do not sum to 1.0. Kernel message: {e}"
        )
    except KernelInvalidInputError as e:
        raise ReconciliationError(f"weighted_score() call failed: {e}")

    weighted_contribution = {d.key: d.score * weights[d.key] for d in reg.dimensions}

    if score >= 4.0:
        raw_verdict = "DEFENSIBLE"
    elif score >= 2.75:
        raw_verdict = "DEFENSIBLE_WITH_MITIGATIONS"
    else:
        raw_verdict = "WEAK_RECOMMEND_COMPETITION"

    label_counts = {"VERIFIED": 0, "ASSERTED": 0, "INFERRED": 0}
    for d in reg.dimensions:
        label_counts[d.evidence_label] += 1
    asserted_majority = label_counts["ASSERTED"] > (len(reg.dimensions) / 2.0)

    verdict_capped = False
    if raw_verdict == "DEFENSIBLE" and asserted_majority:
        verdict = "DEFENSIBLE_WITH_MITIGATIONS"
        verdict_capped = True
    else:
        verdict = raw_verdict

    verdict_color, verdict_bg = VERDICT_COLOR[verdict]

    sorted_dims = sorted(reg.dimensions, key=lambda d: d.score)
    weakest = sorted_dims[0]
    strongest = list(reversed(sorted_dims))[:3]

    effective_mitigations = list(reg.mitigations)
    if verdict_capped:
        effective_mitigations = effective_mitigations + [Mitigation(
            dimension="Multiple (ASSERTED majority)",
            action=(
                "Verify the ASSERTED claims above with current, independently checked evidence. A "
                "majority-ASSERTED case is capped at DEFENSIBLE WITH MITIGATIONS per Hard Rule 2 until "
                "they are independently confirmed."
            ),
            owner="Requester / Compliance",
            due="Before this arrangement is relied on further",
            auto_generated=True,
        )]

    ranked_alternatives = sorted(reg.alternatives, key=_rank_key)
    research_pending = len(reg.research_log) < 2

    if verdict == "WEAK_RECOMMEND_COMPETITION":
        weak_by_mode = {
            "NEW": (
                "run a competitive process (full RFP via rfp-engine, or a lightweight 3-quote comparison, "
                "scaled to value) before award; or re-scope to separate the truly-constrained portion from "
                "the competable portion (Phase 2, dimension 4)"
            ),
            "RENEWAL": (
                "open the next renewal cycle to competition; commission a supplier-landscape shortlist now "
                "so it is ready before the renewal date"
            ),
            "AUDIT": (
                "flag the arrangement for a remediation plan and route to the category lead / Compliance "
                "per sme-matrix.md; this finding is surfaced for the governance owner to act on, not "
                "unilaterally resolved here"
            ),
        }
        text = weak_by_mode[reg.request.mode]
        recommended_next_action = "Recommend: " + text[0].upper() + text[1:] + "."
    elif verdict == "DEFENSIBLE_WITH_MITIGATIONS":
        recommended_next_action = (
            f"Close the {len(effective_mitigations)} mitigation(s) above; the sole-source case is "
            "defensible once they are complete. No competitive process is required this cycle."
        )
    else:
        n_watch = len(effective_mitigations)
        if n_watch:
            recommended_next_action = (
                f"No blocking mitigations. Proceed with the sole-source case; monitor the {n_watch} watch "
                "item(s) noted above at the next review."
            )
        else:
            recommended_next_action = "No blocking mitigations or watch items. Proceed with the sole-source case."

    verdict_sentence = f"Sole-Source Defensibility Score {score:.2f}/5.0: {VERDICT_LABEL[verdict]}."
    if verdict_capped:
        verdict_sentence += (
            " (Capped from a raw Defensible score because a majority of dimensions are labeled ASSERTED "
            "rather than VERIFIED, per Hard Rule 2.)"
        )

    if effective_mitigations:
        matching = [m for m in effective_mitigations if m.dimension.strip().lower() == weakest.label_text.strip().lower()]
        top_mitigation = matching[0] if matching else effective_mitigations[0]
        highest_leverage = (
            f"{top_mitigation.dimension}: {top_mitigation.action} "
            f"(owner: {top_mitigation.owner}, due: {top_mitigation.due})."
        )
    elif verdict == "WEAK_RECOMMEND_COMPETITION" and ranked_alternatives:
        top_alt = ranked_alternatives[0]
        highest_leverage = (
            f"Pursue {top_alt.candidate_name} ({top_alt.reassessed_viability}, confidence "
            f"{top_alt.confidence}) as the leading alternative to re-evaluate."
        )
    else:
        highest_leverage = recommended_next_action

    downstream_map = {
        "DEFENSIBLE": "commercial-negotiation-prep (this verdict is the trigger to move to negotiation strategy)",
        "DEFENSIBLE_WITH_MITIGATIONS": "commercial-negotiation-prep once the mitigations above are closed",
        "WEAK_RECOMMEND_COMPETITION": "supplier-landscape or rfp-engine to run the competitive process this finding recommends",
    }
    downstream = (
        f"Feeds: {downstream_map[verdict]}; executive-summary-package when ready to submit for governance sign-off."
    )

    next_steps = [verdict_sentence, highest_leverage, downstream]

    return GroundTruth(
        weights=weights, weighted_contribution=weighted_contribution, defensibility_score=score,
        raw_verdict=raw_verdict, label_counts=label_counts, asserted_majority=asserted_majority,
        verdict=verdict, verdict_capped=verdict_capped, verdict_label=VERDICT_LABEL[verdict],
        verdict_color=verdict_color, verdict_bg=verdict_bg, weakest=weakest, strongest=strongest,
        effective_mitigations=effective_mitigations, ranked_alternatives=ranked_alternatives,
        research_pending=research_pending, recommended_next_action=recommended_next_action,
        next_steps=next_steps,
    )


# ===========================================================================
# Hard-coded invariant checks (run BEFORE saving)
# ===========================================================================

def _assert_weights_sum_to_one(gt: GroundTruth, tolerance: float = 0.001) -> None:
    total = sum(gt.weights.values())
    if abs(total - 1.0) > tolerance:
        raise ReconciliationError(
            f"Weights-sum-to-one invariant FAILED: the fixed dimension weights sum to {total:.4f}, not "
            f"1.0 (+/- {tolerance}), per SKILL.md Phase 4."
        )


def _assert_weighted_contributions_reconcile(gt: GroundTruth, tolerance: float = 1e-6) -> None:
    summed = sum(gt.weighted_contribution.values())
    if abs(summed - gt.defensibility_score) > tolerance:
        raise ReconciliationError(
            f"Weighted-contributions-reconcile invariant FAILED: the sum of every dimension's weighted "
            f"contribution ({summed}) does not equal the kernel-computed Defensibility Score "
            f"({gt.defensibility_score})."
        )


def _assert_defensibility_score_range(gt: GroundTruth) -> None:
    if not (0.0 <= gt.defensibility_score <= 5.0):
        raise ReconciliationError(
            f"Defensibility-score-range invariant FAILED: score {gt.defensibility_score} is outside [0.0, 5.0]."
        )


def _assert_verdict_matches_score_band_and_cap(gt: GroundTruth) -> None:
    score = gt.defensibility_score
    if score >= 4.0:
        expected_raw = "DEFENSIBLE"
    elif score >= 2.75:
        expected_raw = "DEFENSIBLE_WITH_MITIGATIONS"
    else:
        expected_raw = "WEAK_RECOMMEND_COMPETITION"
    if gt.raw_verdict != expected_raw:
        raise ReconciliationError(
            f"Verdict-matches-score-band invariant FAILED: score {score} implies raw verdict "
            f"{expected_raw!r} (Phase 5 table), got {gt.raw_verdict!r}."
        )
    if gt.raw_verdict == "DEFENSIBLE" and gt.asserted_majority:
        if gt.verdict != "DEFENSIBLE_WITH_MITIGATIONS" or not gt.verdict_capped:
            raise ReconciliationError(
                "Rule-2-cap invariant FAILED: a majority-ASSERTED, raw-Defensible case must be capped to "
                "DEFENSIBLE WITH MITIGATIONS (Hard Rule 2), with verdict_capped True."
            )
    else:
        if gt.verdict != gt.raw_verdict or gt.verdict_capped:
            raise ReconciliationError(
                "Verdict-cap invariant FAILED: verdict should equal raw_verdict, uncapped, when the Rule "
                "2 majority-ASSERTED condition does not apply."
            )


def _assert_price_position_consistent(reg: SoleSourceInput) -> None:
    pv = reg.price_validation
    if pv.method == "none":
        return
    if pv.band_low is None or pv.band_high is None or pv.sole_source_price is None:
        raise ReconciliationError(
            "Price-position-consistent invariant FAILED: method is not 'none' but band_low/band_high/"
            "sole_source_price are not all populated."
        )
    if pv.position == "within" and not (pv.band_low <= pv.sole_source_price <= pv.band_high):
        raise ReconciliationError(
            f"Price-position-consistent invariant FAILED: position is 'within' but sole_source_price "
            f"{pv.sole_source_price} is not inside [{pv.band_low}, {pv.band_high}]."
        )
    if pv.position == "above" and not (pv.sole_source_price > pv.band_high):
        raise ReconciliationError(
            f"Price-position-consistent invariant FAILED: position is 'above' but sole_source_price "
            f"{pv.sole_source_price} is not above {pv.band_high}."
        )
    if pv.position == "below" and not (pv.sole_source_price < pv.band_low):
        raise ReconciliationError(
            f"Price-position-consistent invariant FAILED: position is 'below' but sole_source_price "
            f"{pv.sole_source_price} is not below {pv.band_low}."
        )


def _assert_mitigations_required_when_conditioned(gt: GroundTruth) -> None:
    if gt.verdict != "DEFENSIBLE_WITH_MITIGATIONS":
        return
    if len(gt.effective_mitigations) < 1:
        raise ReconciliationError(
            "Mitigations-required-when-conditioned invariant FAILED: verdict is DEFENSIBLE WITH "
            "MITIGATIONS but no mitigations are present (Phase 5)."
        )
    for m in gt.effective_mitigations:
        if not m.owner.strip() or not m.due.strip():
            raise ReconciliationError(
                f"Mitigations-required-when-conditioned invariant FAILED: mitigation on "
                f"{m.dimension!r} is missing an owner or a due point (Phase 5, 'Each mitigation names "
                "an owner and a due point')."
            )


def _assert_alternatives_nonempty(reg: SoleSourceInput) -> None:
    if len(reg.alternatives) < 1:
        raise ReconciliationError(
            "Alternatives-register-nonempty invariant FAILED: alternatives must contain at least one "
            "row (Deliverable #4)."
        )


def _assert_ranked_alternatives_sorted(gt: GroundTruth) -> None:
    keys = [_rank_key(a) for a in gt.ranked_alternatives]
    if keys != sorted(keys):
        raise ReconciliationError(
            "Ranked-alternatives-sorted invariant FAILED: ranked_alternatives is not sorted by viability "
            "tier, then confidence, then name."
        )


def _assert_gating_items_routed(reg: SoleSourceInput) -> None:
    haystack = " ".join(
        [reg.request.category, reg.request.need_description, reg.request.stated_rationale]
        + [d.evidence + " " + d.rationale for d in reg.dimensions]
    ).lower()
    hits = [kw for kw in GATING_KEYWORDS if kw in haystack]
    if hits and not reg.sme_routing:
        raise ReconciliationError(
            f"Gating-items-routed invariant FAILED: the input text references {hits}, a potential "
            "debarment/sanctions/GxP-adjacent gating concern (Hard Rule 6, 'Gating items are flagged, "
            "never cleared'), but sme_routing is empty. Route it, do not silently clear it."
        )


def run_hardcoded_invariant_checks(reg: SoleSourceInput, gt: GroundTruth) -> None:
    """Run every hard-coded invariant. Raises ReconciliationError on any
    failure; callers must not save the document if this raises."""
    _assert_weights_sum_to_one(gt)
    _assert_weighted_contributions_reconcile(gt)
    _assert_defensibility_score_range(gt)
    _assert_verdict_matches_score_band_and_cap(gt)
    _assert_price_position_consistent(reg)
    _assert_mitigations_required_when_conditioned(gt)
    _assert_alternatives_nonempty(reg)
    _assert_ranked_alternatives_sorted(gt)
    _assert_gating_items_routed(reg)


# ===========================================================================
# Shared formatting helpers (format-agnostic; used by both builders)
# ===========================================================================

def _usd_or_na(v: Optional[float]) -> str:
    return "N/A" if v is None else f"${v:,.0f}"


def _pct(frac: float) -> str:
    return f"{frac * 100:.0f}%"


def _est_value_text(reg: SoleSourceInput) -> str:
    return "NEEDS_INPUT, no estimated value provided" if reg.request.est_value_usd is None else _usd_or_na(reg.request.est_value_usd)


def _price_position_sentence(reg: SoleSourceInput) -> str:
    pv = reg.price_validation
    if pv.method == "none":
        return "No price validation was attempted; the sole-source price is taken on faith (Rule 5)."
    pos_text = {
        "within": "within the validated band, independently supporting the sole-source price",
        "above": "above the validated band, a flag worth raising in negotiation",
        "below": "below the validated band",
    }.get(pv.position, pv.position)
    return (
        f"The sole-source price ({_usd_or_na(pv.sole_source_price)}) sits {pos_text} "
        f"({_usd_or_na(pv.band_low)} to {_usd_or_na(pv.band_high)}, method: {pv.method}, source: {pv.source})."
    )


def _research_pending_sentence(reg: SoleSourceInput, gt: GroundTruth) -> str:
    n = len(reg.research_log)
    if gt.research_pending:
        return (
            f"RESEARCH PENDING: only {n} research-log entr{'y' if n == 1 else 'ies'} on file; the G7 "
            "minimum (2 independent searches) was not met. Treat Alternative Supplier Availability "
            "evidence as thin until a fuller check (ideally a full supplier-landscape run) is completed."
        )
    return f"G7 minimum met: {n} research-log entries on file (minimum 2)."


# ===========================================================================
# 2a. DOCX builder (python-docx)
# ===========================================================================

def _require_docx() -> None:
    if not DOCX_AVAILABLE:
        raise ImportError(
            "python-docx is not installed in this Python environment, so no .docx file can be written. "
            "Install it (`pip install python-docx`) or point this script at an interpreter that already "
            f"has it. Original import error: {_DOCX_IMPORT_ERROR}"
        )


def _docx_set_cell_background(cell, hex_color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(docx_qn("w:val"), "clear")
    shd.set(docx_qn("w:color"), "auto")
    shd.set(docx_qn("w:fill"), hex_color)
    tc_pr.append(shd)


def _docx_set_cell_text(cell, text: str, bold: bool = False, color_hex: str = LILLY_BLACK,
                         size_pt: int = 9, align=None) -> None:
    cell.text = str(text)
    para = cell.paragraphs[0]
    if align is not None:
        para.alignment = align
    run = para.runs[0] if para.runs else para.add_run("")
    run.font.bold = bold
    run.font.size = DocxPt(size_pt)
    run.font.name = "Calibri"
    run.font.color.rgb = DocxRGBColor.from_string(color_hex)


def _docx_add_heading(doc, text: str, level: int) -> None:
    spec = {1: (14, BOLD_BLUE, 18, 8), 2: (12, LILLY_BLACK, 14, 6), 3: (11, LILLY_BLACK, 10, 4)}
    size_pt, color_hex, before_pt, after_pt = spec.get(level, spec[3])
    para = doc.add_paragraph()
    para.paragraph_format.space_before = DocxPt(before_pt)
    para.paragraph_format.space_after = DocxPt(after_pt)
    run = para.add_run(text)
    run.font.bold = True
    run.font.size = DocxPt(size_pt)
    run.font.name = "Calibri"
    run.font.color.rgb = DocxRGBColor.from_string(color_hex)


def _docx_add_body(doc, text: str, bold: bool = False, italic: bool = False,
                    size_pt: int = 10, color_hex: str = LILLY_BLACK) -> None:
    para = doc.add_paragraph()
    para.paragraph_format.space_after = DocxPt(6)
    run = para.add_run(text)
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = DocxPt(size_pt)
    run.font.name = "Calibri"
    run.font.color.rgb = DocxRGBColor.from_string(color_hex)


def _docx_add_bullets(doc, items: Sequence[str]) -> None:
    for item in items:
        para = doc.add_paragraph(style="List Bullet")
        run = para.add_run(item)
        run.font.size = DocxPt(10)
        run.font.name = "Calibri"
        run.font.color.rgb = DocxRGBColor.from_string(LILLY_BLACK)


def _docx_add_table(doc, headers: Sequence[str], rows: Sequence[Sequence[Any]]):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        _docx_set_cell_background(hdr_cells[i], LILLY_RED)
        _docx_set_cell_text(hdr_cells[i], h, bold=True, color_hex=WHITE, size_pt=9)
    for row_values in rows:
        cells = table.add_row().cells
        for i, val in enumerate(row_values):
            _docx_set_cell_text(cells[i], str(val), size_pt=9)
    return table


def build_docx(reg: SoleSourceInput, gt: GroundTruth):
    """Build the sole_source_challenge_report.docx per SKILL.md Deliverable
    #1: Sections 01-06 plus an Appendix, the same skeleton every run. This
    function does NOT save the file and does NOT run the reconciliation
    checks; call generate_sole_source_challenge() for the full validated
    pipeline."""
    _require_docx()
    doc = Document()

    footer_para = doc.sections[0].footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_para.add_run(
        "CONFIDENTIAL, internal use only. Lilly Procurement Skills Suite, Sole-Source Challenge."
    )
    footer_run.font.size = DocxPt(8)
    footer_run.font.name = "Calibri"
    footer_run.font.color.rgb = DocxRGBColor.from_string(MUTED_GREY)

    # --- Title page ------------------------------------------------------
    title_para = doc.add_paragraph()
    title_para.paragraph_format.space_after = DocxPt(4)
    title_run = title_para.add_run("Sole-Source Challenge and Justification")
    title_run.font.bold = True
    title_run.font.size = DocxPt(26)
    title_run.font.name = "Calibri"
    title_run.font.color.rgb = DocxRGBColor.from_string(LILLY_BLACK)

    scope_para = doc.add_paragraph()
    scope_para.paragraph_format.space_after = DocxPt(10)
    scope_run = scope_para.add_run(f"{reg.request.need_description}")
    scope_run.font.size = DocxPt(11)
    scope_run.font.name = "Calibri"
    scope_run.font.color.rgb = DocxRGBColor.from_string(BOLD_BLUE)

    conf_para = doc.add_paragraph()
    conf_para.paragraph_format.space_after = DocxPt(10)
    conf_run = conf_para.add_run("CONFIDENTIAL, INTERNAL USE ONLY")
    conf_run.font.bold = True
    conf_run.font.size = DocxPt(10)
    conf_run.font.name = "Calibri"
    conf_run.font.color.rgb = DocxRGBColor.from_string(LILLY_RED)

    for line in [
        f"Proposed / incumbent supplier: {reg.request.supplier}",
        f"Mode: {reg.request.mode}",
        f"Category: {reg.request.category}",
        f"Estimated value: {_est_value_text(reg)}",
        f"Date: {reg.request.date}",
    ]:
        _docx_add_body(doc, line)
    doc.add_page_break()

    # --- Section 01: Request Summary --------------------------------------
    _docx_add_heading(doc, "01 Request Summary", level=1)
    _docx_add_body(doc, f"Need / scope: {reg.request.need_description}", bold=True)
    _docx_add_body(doc, f"Proposed / incumbent supplier: {reg.request.supplier}")
    _docx_add_body(doc, f"Mode: {reg.request.mode}, {MODE_DEFINITION[reg.request.mode]}")
    _docx_add_body(doc, f"Requester's stated rationale: \"{reg.request.stated_rationale}\"", italic=True)
    _docx_add_table(
        doc, ["Field", "Value"],
        [
            ["Category", reg.request.category],
            ["Estimated value", _est_value_text(reg)],
            ["Term", reg.request.term],
            ["Requester", reg.request.requester],
            ["As-of date", reg.request.date],
        ],
    )

    # --- Section 02: The Challenge -----------------------------------------
    _docx_add_heading(doc, "02 The Challenge", level=1)
    _docx_add_body(
        doc,
        "Seven fixed dimensions (SKILL.md Phase 4), each scored 0.0-5.0 against a defensible-vs-weak "
        "anchor, each claim labeled VERIFIED, ASSERTED, or INFERRED (Hard Rule 2).",
    )
    _docx_add_table(
        doc, ["#", "Dimension", "Weight", "Score", "Label", "Confidence"],
        [
            [d.order, e.label_text, _pct(e.weight), f"{e.score:.1f}", e.evidence_label, e.confidence]
            for d, e in zip(DIMENSION_SPECS, reg.dimensions)
        ],
    )
    _docx_add_heading(doc, "Evidence and rationale, by dimension", level=2)
    _docx_add_bullets(doc, [
        f"{e.label_text} ({e.evidence_label}, {e.confidence} confidence): {e.rationale} Evidence: {e.evidence}"
        for e in reg.dimensions
    ])
    _docx_add_body(
        doc,
        f"Label distribution: {gt.label_counts['VERIFIED']} VERIFIED, {gt.label_counts['ASSERTED']} "
        f"ASSERTED, {gt.label_counts['INFERRED']} INFERRED."
        + (" Majority ASSERTED: the Hard Rule 2 verdict cap applies (see Section 04)." if gt.asserted_majority else ""),
    )

    # --- Section 03: Evidence and Price Validation --------------------------
    _docx_add_heading(doc, "03 Evidence and Price Validation", level=1)
    _docx_add_heading(doc, "Research log", level=2)
    if reg.research_log:
        _docx_add_table(
            doc, ["Query", "Source", "Date", "Results"],
            [[r.query, r.source, r.date, r.results] for r in reg.research_log],
        )
    else:
        _docx_add_body(doc, "NEEDS_INPUT, no research log entries on file for this run.")
    _docx_add_body(doc, _research_pending_sentence(reg, gt))
    _docx_add_heading(doc, "Consumed artifacts", level=2)
    if reg.consumed_artifacts:
        _docx_add_bullets(doc, reg.consumed_artifacts)
    else:
        _docx_add_body(doc, "No upstream artifacts (supplier-landscape, market-rate-benchmarking, should-cost-builder) were consumed this run.")
    _docx_add_heading(doc, "Price validation", level=2)
    pv = reg.price_validation
    _docx_add_table(
        doc, ["Method", "Source", "Band Low", "Band High", "Sole-Source Price", "Position"],
        [[pv.method, pv.source or "N/A", _usd_or_na(pv.band_low), _usd_or_na(pv.band_high),
          _usd_or_na(pv.sole_source_price), pv.position]],
    )
    _docx_add_body(doc, _price_position_sentence(reg))

    # --- Section 04: Verdict and Recommendation -----------------------------
    _docx_add_heading(doc, "04 Verdict and Recommendation", level=1)
    _docx_add_body(doc, gt.next_steps[0], bold=True, color_hex=gt.verdict_color)
    _docx_add_body(
        doc,
        f"Calc: sum over 7 dimensions of (score x weight), computed by numeric_kernel.weighted_score() "
        f"= {gt.defensibility_score:.3f} / 5.0.",
    )
    _docx_add_body(
        doc,
        "This is a recommendation, not an approval (Hard Rule 3). Route the actual threshold / "
        "approval-chain question to process-navigator.",
        italic=True,
    )

    if gt.verdict == "DEFENSIBLE":
        _docx_add_heading(doc, "Leading dimensions", level=2)
        _docx_add_bullets(doc, [f"{d.label_text}: {d.score:.1f}/5.0, {d.rationale}" for d in gt.strongest])
        _docx_add_heading(doc, "Watch items", level=2)
        if gt.effective_mitigations:
            _docx_add_table(
                doc, ["Dimension", "Watch item", "Owner", "Due"],
                [[m.dimension, m.action, m.owner, m.due] for m in gt.effective_mitigations],
            )
        else:
            _docx_add_body(doc, "No watch items identified this run.")
    elif gt.verdict == "DEFENSIBLE_WITH_MITIGATIONS":
        _docx_add_heading(doc, "Mitigations (condition the justification on closing these)", level=2)
        _docx_add_table(
            doc, ["Dimension", "Action", "Owner", "Due"],
            [[m.dimension, m.action, m.owner, m.due] for m in gt.effective_mitigations],
        )
    else:
        _docx_add_heading(doc, "Ranked alternatives worth pursuing", level=2)
        _docx_add_table(
            doc, ["Rank", "Candidate", "Origin", "Viability", "Confidence", "Capability gap"],
            [
                [i, a.candidate_name, a.origin, a.reassessed_viability, a.confidence, a.capability_gap]
                for i, a in enumerate(gt.ranked_alternatives, start=1)
            ],
        )

    _docx_add_heading(doc, "Alternatives considered", level=2)
    _docx_add_table(
        doc, ["Candidate", "Origin", "Original exclusion reason", "Capability gap", "Reassessed viability", "Confidence", "Source", "Date"],
        [
            [a.candidate_name, a.origin, a.original_exclusion_reason, a.capability_gap,
             a.reassessed_viability, a.confidence, a.source, a.date]
            for a in reg.alternatives
        ],
    )
    _docx_add_heading(doc, "Recommended next action", level=2)
    _docx_add_body(doc, gt.recommended_next_action, bold=True)

    # --- Section 05: Next Steps and SME Routing -----------------------------
    _docx_add_heading(doc, "05 Next Steps and SME Routing", level=1)
    _docx_add_bullets(doc, gt.next_steps)
    _docx_add_heading(doc, "SME routing", level=2)
    if reg.sme_routing:
        _docx_add_table(
            doc, ["Issue", "Route to", "Reason"],
            [[s.issue, s.route_to, s.reason] for s in reg.sme_routing],
        )
    else:
        _docx_add_body(doc, "No gating item was surfaced this run; no SME routing required.")

    # --- Section 06: Research Methodology -----------------------------------
    _docx_add_heading(doc, "06 Research Methodology", level=1)
    _docx_add_body(doc, reg.research_methodology_note)
    _docx_add_body(doc, _research_pending_sentence(reg, gt))
    _docx_add_body(
        doc,
        f"Label distribution across the 7 dimensions: {gt.label_counts['VERIFIED']} VERIFIED, "
        f"{gt.label_counts['ASSERTED']} ASSERTED, {gt.label_counts['INFERRED']} INFERRED.",
    )
    _docx_add_body(doc, f"Mode framing: {reg.request.mode}, {MODE_DEFINITION[reg.request.mode]}")

    # --- Appendix: Raw Scorecard --------------------------------------------
    _docx_add_heading(doc, "Appendix: Raw Scorecard", level=1)
    rows = [
        [e.label_text, _pct(e.weight), f"{e.score:.1f}", f"{gt.weighted_contribution[e.key]:.3f}",
         e.evidence_label, e.confidence, e.evidence]
        for e in reg.dimensions
    ]
    rows.append(["TOTAL (Defensibility Score)", "100%", "", f"{gt.defensibility_score:.3f}", "", "", ""])
    _docx_add_table(doc, ["Dimension", "Weight", "Score", "Weighted contribution", "Label", "Confidence", "Evidence"], rows)
    _docx_add_body(
        doc,
        "REFLECT-ONLY: this skill reads and drafts; it never sends, submits, files, or updates anything "
        "in ARIA, Ariba, LEAH, Aravo, ServiceNow, SAP, or any other Lilly system of record.",
        italic=True,
    )

    return doc


# ===========================================================================
# 2b. PPTX builder (python-pptx). One idea per slide: title plus concise
# bullets/table. Same section order as the DOCX skeleton (01-06, Appendix).
# ===========================================================================

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

def _require_pptx() -> None:
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is not installed in this Python environment, so no .pptx file can be written. "
            "Install it (`pip install python-pptx`) or point this script at an interpreter that already "
            f"has it. Original import error: {_PPTX_IMPORT_ERROR}"
        )


def _pptx_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() == "blank":
            return layout
    return prs.slide_layouts[6]


def _pptx_new_slide(prs):
    return prs.slides.add_slide(_pptx_blank_layout(prs))


def _pptx_rect(slide, left_in, top_in, width_in, height_in, fill_hex, line_hex=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = PptxRGBColor.from_string(fill_hex)
    if line_hex:
        shp.line.color.rgb = PptxRGBColor.from_string(line_hex)
        shp.line.width = PptxPt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _pptx_textbox(slide, left_in, top_in, width_in, height_in, text, size=12, bold=False,
                   color_hex=LILLY_BLACK, align=None, font="Calibri", italic=False):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = PptxRGBColor.from_string(color_hex)
    return box


def _pptx_bullets(slide, items, left_in, top_in, width_in, height_in, size=13):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = PptxPt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = PptxRGBColor.from_string(LILLY_BLACK)
        p.space_after = PptxPt(8)
    return box


def _pptx_table(slide, headers, rows, left_in, top_in, width_in, height_in, body_size=10.5):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptxRGBColor.from_string(LILLY_RED)
        para = cell.text_frame.paragraphs[0]
        para.font.size = PptxPt(11)
        para.font.bold = True
        para.font.color.rgb = PptxRGBColor.from_string(WHITE)
        para.font.name = "Calibri"
    for r, row_vals in enumerate(rows, start=1):
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptxRGBColor.from_string(WHITE if r % 2 else NEUTRAL_STONE)
            para = cell.text_frame.paragraphs[0]
            para.font.size = PptxPt(body_size)
            para.font.name = "Calibri"
            para.font.color.rgb = PptxRGBColor.from_string(LILLY_BLACK)
    return table


def _pptx_header(slide, kicker, title):
    _pptx_rect(slide, 0, 0, SLIDE_W_IN, 1.0, LILLY_BLACK)
    _pptx_textbox(slide, 0.5, 0.12, 10.0, 0.3, kicker, size=11, bold=True, color_hex=LILLY_RED)
    _pptx_textbox(slide, 0.5, 0.42, 12.0, 0.55, title, size=22, bold=True, color_hex=WHITE)
    _pptx_textbox(
        slide, 0.5, SLIDE_H_IN - 0.35, 10.0, 0.3,
        "CONFIDENTIAL, internal use only. Lilly Procurement Skills Suite, Sole-Source Challenge.",
        size=8, color_hex=MUTED_GREY,
    )


def build_pptx(reg: SoleSourceInput, gt: GroundTruth):
    """Build the sole_source_challenge_report.pptx: one idea per slide,
    mirroring the DOCX skeleton's section order (01-06, Appendix). Does NOT
    save the file and does NOT run reconciliation checks; call
    generate_sole_source_challenge() for the full validated pipeline."""
    _require_pptx()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # --- Slide 1: Title -----------------------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, WHITE)
    _pptx_rect(slide, 0, 0, SLIDE_W_IN, 0.15, LILLY_RED)
    _pptx_textbox(slide, 0.8, 2.0, 11.7, 1.2, "Sole-Source Challenge and Justification", size=36, bold=True, color_hex=LILLY_BLACK)
    _pptx_textbox(slide, 0.8, 3.1, 11.7, 0.9, reg.request.need_description, size=16, color_hex=BOLD_BLUE)
    _pptx_textbox(
        slide, 0.8, 3.9, 11.7, 0.4,
        f"Supplier: {reg.request.supplier}    Mode: {reg.request.mode}    Date: {reg.request.date}",
        size=13, color_hex=LILLY_BLACK,
    )
    _pptx_textbox(slide, 0.8, 4.4, 6.0, 0.4, "CONFIDENTIAL, INTERNAL USE ONLY", size=12, bold=True, color_hex=LILLY_RED)

    # --- Slide 2: 01 Request Summary -----------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "01", "Request Summary")
    _pptx_bullets(slide, [
        f"Need / scope: {reg.request.need_description}",
        f"Mode: {reg.request.mode}, {MODE_DEFINITION[reg.request.mode]}",
        f"Requester's stated rationale: \"{reg.request.stated_rationale}\"",
    ], 0.5, 1.2, 12.3, 2.1, size=14)
    _pptx_table(
        slide, ["Field", "Value"],
        [
            ["Supplier", reg.request.supplier],
            ["Category", reg.request.category],
            ["Estimated value", _est_value_text(reg)],
            ["Term", reg.request.term],
            ["Requester", reg.request.requester],
        ],
        0.5, 3.4, 12.3, 3.2,
    )

    # --- Slide 3: 02 The Challenge, Scorecard --------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "02", "The Challenge, Scorecard")
    _pptx_table(
        slide, ["#", "Dimension", "Weight", "Score /5.0", "Label", "Confidence"],
        [
            [d.order, e.label_text, _pct(e.weight), f"{e.score:.1f}", e.evidence_label, e.confidence]
            for d, e in zip(DIMENSION_SPECS, reg.dimensions)
        ],
        0.5, 1.2, 12.3, 5.4,
    )

    # --- Slide 4: 02 The Challenge, Evidence and Rationale -------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "02", "The Challenge, Evidence and Rationale")
    _pptx_bullets(slide, [
        f"{e.label_text} ({e.evidence_label}): {e.rationale[:180]}"
        for e in reg.dimensions
    ], 0.5, 1.2, 12.3, 5.5, size=12)

    # --- Slide 5: 03 Evidence, Research Log ----------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "03", "Evidence, Research Log")
    if reg.research_log:
        _pptx_table(
            slide, ["Query", "Source", "Date", "Results"],
            [[r.query, r.source, r.date, r.results] for r in reg.research_log],
            0.5, 1.2, 12.3, 3.6,
        )
    else:
        _pptx_textbox(slide, 0.5, 1.3, 12.0, 0.5, "NEEDS_INPUT, no research log entries on file for this run.", size=14)
    _pptx_textbox(slide, 0.5, 5.1, 12.3, 0.8, _research_pending_sentence(reg, gt), size=13, color_hex=AMBER if gt.research_pending else LILLY_BLACK)
    consumed_text = "; ".join(reg.consumed_artifacts) if reg.consumed_artifacts else "No upstream artifacts consumed this run."
    _pptx_textbox(slide, 0.5, 5.8, 12.3, 0.9, f"Consumed artifacts: {consumed_text}", size=12, italic=True)

    # --- Slide 6: 03 Price Validation ----------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "03", "Price Validation")
    pv = reg.price_validation
    _pptx_table(
        slide, ["Method", "Source", "Band Low", "Band High", "Sole-Source Price", "Position"],
        [[pv.method, pv.source or "N/A", _usd_or_na(pv.band_low), _usd_or_na(pv.band_high),
          _usd_or_na(pv.sole_source_price), pv.position]],
        0.5, 1.4, 12.3, 1.2,
    )
    _pptx_textbox(slide, 0.5, 3.0, 12.3, 1.5, _price_position_sentence(reg), size=15)

    # --- Slide 7: 04 Verdict --------------------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "04", "Verdict")
    _pptx_rect(slide, 0.5, 1.4, 5.5, 2.6, gt.verdict_bg, line_hex=gt.verdict_color)
    _pptx_textbox(slide, 0.8, 1.65, 5.0, 1.2, f"{gt.defensibility_score:.2f}", size=54, bold=True, color_hex=gt.verdict_color)
    _pptx_textbox(slide, 0.8, 3.05, 5.0, 0.8, gt.verdict_label, size=18, bold=True, color_hex=LILLY_BLACK)
    _pptx_textbox(
        slide, 6.3, 1.4, 6.5, 2.6,
        f"Calc: sum over 7 dimensions of (score x weight), computed by numeric_kernel.weighted_score() "
        f"= {gt.defensibility_score:.3f} / 5.0."
        + (" Capped from a raw Defensible score, majority ASSERTED per Hard Rule 2." if gt.verdict_capped else ""),
        size=13,
    )
    _pptx_textbox(
        slide, 0.5, 4.3, 12.3, 1.0,
        "This is a recommendation, not an approval (Hard Rule 3). Route the actual threshold / "
        "approval-chain question to process-navigator.",
        size=12, italic=True,
    )

    # --- Slide 8: 04 Mitigations / Watch Items --------------------------------
    slide = _pptx_new_slide(prs)
    if gt.verdict == "DEFENSIBLE_WITH_MITIGATIONS":
        _pptx_header(slide, "04", "Mitigations")
        _pptx_table(
            slide, ["Dimension", "Action", "Owner", "Due"],
            [[m.dimension, m.action, m.owner, m.due] for m in gt.effective_mitigations],
            0.5, 1.2, 12.3, 5.4,
        )
    elif gt.verdict == "DEFENSIBLE":
        _pptx_header(slide, "04", "Leading Dimensions and Watch Items")
        _pptx_bullets(slide, [f"{d.label_text}: {d.score:.1f}/5.0, {d.rationale[:140]}" for d in gt.strongest], 0.5, 1.2, 12.3, 2.4, size=13)
        if gt.effective_mitigations:
            _pptx_table(
                slide, ["Dimension", "Watch item", "Owner", "Due"],
                [[m.dimension, m.action, m.owner, m.due] for m in gt.effective_mitigations],
                0.5, 3.8, 12.3, 2.8,
            )
        else:
            _pptx_textbox(slide, 0.5, 3.9, 12.0, 0.5, "No watch items identified this run.", size=14)
    else:
        _pptx_header(slide, "04", "Ranked Alternatives Worth Pursuing")
        _pptx_table(
            slide, ["Rank", "Candidate", "Origin", "Viability", "Confidence", "Capability gap"],
            [
                [i, a.candidate_name, a.origin, a.reassessed_viability, a.confidence, a.capability_gap]
                for i, a in enumerate(gt.ranked_alternatives, start=1)
            ],
            0.5, 1.2, 12.3, 5.4,
        )

    # --- Slide 9: 04 Alternatives Considered -----------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "04", "Alternatives Considered")
    _pptx_table(
        slide, ["Candidate", "Origin", "Viability", "Confidence", "Capability gap"],
        [
            [a.candidate_name, a.origin, a.reassessed_viability, a.confidence, a.capability_gap]
            for a in reg.alternatives
        ],
        0.5, 1.2, 12.3, 4.7,
    )
    _pptx_textbox(slide, 0.5, 6.05, 12.3, 0.9, f"Recommended next action: {gt.recommended_next_action}", size=13, bold=True)

    # --- Slide 10: 05 Next Steps ------------------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "05", "Next Steps")
    _pptx_bullets(slide, gt.next_steps, 0.5, 1.3, 12.3, 4.5, size=16)

    # --- Slide 11: 05 SME Routing ------------------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "05", "SME Routing")
    if reg.sme_routing:
        _pptx_table(
            slide, ["Issue", "Route to", "Reason"],
            [[s.issue, s.route_to, s.reason] for s in reg.sme_routing],
            0.5, 1.3, 12.3, 4.0,
        )
    else:
        _pptx_textbox(slide, 0.5, 1.4, 12.0, 0.6, "No gating item was surfaced this run; no SME routing required.", size=15)

    # --- Slide 12: 06 Research Methodology ---------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "06", "Research Methodology")
    _pptx_bullets(slide, [
        reg.research_methodology_note,
        _research_pending_sentence(reg, gt),
        f"Label distribution: {gt.label_counts['VERIFIED']} VERIFIED, {gt.label_counts['ASSERTED']} ASSERTED, {gt.label_counts['INFERRED']} INFERRED.",
        f"Mode framing: {reg.request.mode}, {MODE_DEFINITION[reg.request.mode]}",
    ], 0.5, 1.3, 12.3, 5.2, size=14)

    # --- Slide 13: Appendix, Raw Scorecard ----------------------------------------
    slide = _pptx_new_slide(prs)
    _pptx_header(slide, "Appendix", "Raw Scorecard")
    rows = [
        [e.label_text, _pct(e.weight), f"{e.score:.1f}", f"{gt.weighted_contribution[e.key]:.3f}", e.evidence_label]
        for e in reg.dimensions
    ]
    rows.append(["TOTAL (Defensibility Score)", "100%", "", f"{gt.defensibility_score:.3f}", ""])
    _pptx_table(slide, ["Dimension", "Weight", "Score", "Weighted contribution", "Label"], rows, 0.5, 1.2, 12.3, 5.4)

    return prs


# ===========================================================================
# Full pipeline: validate -> compute ground truth -> hard-invariant checks
# -> build document -> save
# ===========================================================================

def generate_sole_source_challenge(raw_register: Dict[str, Any], output_path: str,
                                    fmt: str) -> Tuple[SoleSourceInput, GroundTruth]:
    """End-to-end: validate the raw sole-source register, compute ground
    truth via numeric_kernel, run the hard-coded invariant checks, build
    the document in the requested format ('pptx' or 'docx'), and only then
    save it. Raises rather than saving a document that fails validation or
    reconciliation."""
    fmt_norm = str(fmt).strip().lower()
    if fmt_norm not in ("pptx", "docx"):
        raise ValueError(f"fmt must be 'pptx' or 'docx'; got {fmt!r}.")
    register = copy.deepcopy(raw_register)
    reg = validate_sole_source_input(register)
    gt = compute_ground_truth(reg)
    run_hardcoded_invariant_checks(reg, gt)
    if fmt_norm == "docx":
        doc = build_docx(reg, gt)
        doc.save(output_path)
    else:
        prs = build_pptx(reg, gt)
        prs.save(output_path)
    return reg, gt


# ===========================================================================
# Demo data (clearly illustrative; no fabrication of real Lilly facts)
# ===========================================================================

def _demo_defensible_with_mitigations_register() -> Dict[str, Any]:
    """ILLUSTRATIVE DEMO DATA, not a real sourcing event. Adapted directly
    from SKILL.md Phase 4's own worked numeric example (the seven dimension
    scores below reproduce that example exactly, expected Defensibility
    Score 3.475) and from the illustrative GxP Environmental Monitoring
    System renewal scenario in this skill's reference dashboard,
    examples/sole_source_challenge_canonical_dashboard.jsx, so this
    generator, the SKILL.md prose, and the canonical dashboard all
    reconcile on the same numbers. RENEWAL mode; majority VERIFIED, so no
    Hard Rule 2 cap fires (raw and final verdict are both DEFENSIBLE WITH
    MITIGATIONS)."""
    return {
        "request": {
            "supplier": "Vaisala viewLinc (incumbent)",
            "need_description": "GxP Environmental Monitoring System (EMS), Renewal, 3 Manufacturing Sites (illustrative demo scenario)",
            "stated_rationale": "It is the incumbent system; Site Quality says no alternative has been evaluated since the original 2019 installation.",
            "mode": "RENEWAL",
            "date": "July 22, 2026",
            "category": "Lab and Clinical, Validated Instrumentation / GxP Software (illustrative)",
            "requester": "Site Quality and Validation, Indianapolis (illustrative)",
            "est_value_usd": 1240000,
            "term": "3 years, annual license plus calibration services",
        },
        "dimensions": [
            {"dimension_key": "unique_capability", "score": 4.5, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "Validation master plan VMP-EMS-2024, Site Quality file (uploaded)",
             "rationale": "viewLinc is the only continuously-monitoring EMS with completed IQ/OQ/PQ validation across all 3 controlled environments; swapping platforms would trigger a multi-month revalidation before any site could go live again."},
            {"dimension_key": "constraint_basis", "score": 4.0, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "21 CFR Part 11 qualification record, Site Quality file (uploaded)",
             "rationale": "Safety/regulatory qualification: the existing sensor network and 21 CFR Part 11 audit trail are already qualified against this specific instrument base; re-qualifying a different vendor's hardware is a structural, not preferential, constraint."},
            {"dimension_key": "competition_history", "score": 2.0, "label": "VERIFIED", "confidence": "MEDIUM",
             "evidence": "No RFP/RFI on file for this renewal cycle (M365 search, this run)",
             "rationale": "No formal market scan was conducted at this renewal. The last competitive review was at original installation in 2019, over 6 years ago; the requester's claim that nothing else was considered is accurate but weak on its own."},
            {"dimension_key": "requirements_separability", "score": 3.5, "label": "VERIFIED", "confidence": "MEDIUM",
             "evidence": "System architecture diagram, EMS-ARCH-2024 (uploaded)",
             "rationale": "Core continuous monitoring and alarming cannot be split from the validated instrumentation without breaking the qualification. The reporting and analytics layer, however, is architecturally separable and could in principle be competitively sourced at the next platform upgrade."},
            {"dimension_key": "alt_availability", "score": 2.5, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "Light market-check research log, this run (2 searches, RESEARCH PENDING for a full shortlist)",
             "rationale": "A light 2-search market-check surfaced two plausible alternative EMS platforms. Neither has been seriously evaluated against the validated instrument base; a full supplier-landscape run would firm this up before the next renewal."},
            {"dimension_key": "urgency_legitimacy", "score": 3.0, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "License term end date, existing MSA exhibit (uploaded)",
             "rationale": "The renewal date is fixed by the existing license term (external), but the request arrived inside the standard renewal window with adequate lead time; this is a mix, not a hard external emergency."},
            {"dimension_key": "price_validation", "score": 4.0, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "should_cost_model.xlsx (should-cost-builder output, consumed this run)",
             "rationale": "A should-cost-builder range exists for the annual license plus calibration-services fee. The sole-source ask sits within the modeled range, independently validating the price despite no competitive bid."},
        ],
        "mitigations": [
            {"dimension": "Competition History",
             "action": "Document a formal market scan before the next renewal cycle, not just an informal check.",
             "owner": "Site Quality and Validation lead", "due": "Before next renewal (24 months)"},
            {"dimension": "Alternative Supplier Availability",
             "action": "Formalize the 2 alternative EMS platforms surfaced in this run's light market-check as a standing comparison set; consider a full supplier-landscape run given the contract value.",
             "owner": "Requester", "due": "90 days"},
            {"dimension": "Requirements Separability",
             "action": "Evaluate decoupling the reporting and analytics layer from the core monitoring at the next platform upgrade so at least part of the scope can be competed.",
             "owner": "IT Procurement category lead", "due": "Next major upgrade cycle"},
        ],
        "alternatives": [
            {"candidate_name": "Rees Scientific", "origin": "market-check",
             "original_exclusion_reason": "n/a, newly surfaced this run",
             "capability_gap": "No existing validation on Lilly's controlled environments; would require full IQ/OQ/PQ revalidation across 3 sites.",
             "reassessed_viability": "Viable With Gaps", "confidence": "MEDIUM", "source": "Web market-check, this run", "date": "Jul 2026"},
            {"candidate_name": "Veriteq / GE Vionic", "origin": "market-check",
             "original_exclusion_reason": "n/a, newly surfaced this run",
             "capability_gap": "Comparable sensor network but no current Lilly site integration; migration cost not modeled.",
             "reassessed_viability": "Viable With Gaps", "confidence": "LOW", "source": "Web market-check, this run", "date": "Jul 2026"},
            {"candidate_name": "ClimateCheck Systems", "origin": "supplier-landscape-excluded",
             "original_exclusion_reason": "insufficient_evidence: no disclosed GxP validation track record (excluded_vendors.csv)",
             "capability_gap": "No public evidence of pharma-grade qualification support.",
             "reassessed_viability": "Not Viable", "confidence": "HIGH",
             "source": "supplier-landscape excluded_vendors.csv (consumed this run)", "date": "Jun 2026"},
        ],
        "price_validation": {
            "method": "should-cost",
            "source": "should_cost_model.xlsx (should-cost-builder, annual license plus calibration services)",
            "band_low": 380000, "band_high": 445000, "sole_source_price": 402000, "position": "within",
        },
        "research_log": [
            {"query": "GxP environmental monitoring system alternative to Vaisala viewLinc", "source": "Web search", "date": "Jul 21, 2026", "results": 6},
            {"query": "continuous temperature monitoring pharma manufacturing vendor comparison 2026", "source": "Web search", "date": "Jul 21, 2026", "results": 5},
            {"query": "sole-source renewal history, EMS platform", "source": "M365 (SharePoint, this tenant)", "date": "Jul 21, 2026", "results": 0},
        ],
        "consumed_artifacts": [
            "excluded_vendors.csv (supplier-landscape, consumed for ClimateCheck Systems exclusion)",
            "should_cost_model.xlsx (should-cost-builder, price validation)",
        ],
        "sme_routing": [
            {"issue": "GxP qualification status of any alternative EMS platform", "route_to": "Quality / GxP",
             "reason": "Any alternative must clear a validation review before it can be treated as viable; this skill flags, Quality decides."},
            {"issue": "Whether this renewal value re-triggers FRAP review", "route_to": "process-navigator",
             "reason": "Threshold and system-requirement questions route to process-navigator's live policy read, not asserted here."},
        ],
        "research_methodology_note": (
            "Light market-check (G7 minimum, 2 independent web searches) plus one M365 search for a "
            "prior competitive event; see the Research Log below. A full supplier-landscape run was not "
            "commissioned this cycle."
        ),
    }


def _demo_weak_register() -> Dict[str, Any]:
    """ILLUSTRATIVE DEMO DATA, a fictional NEW-mode marketing services
    scenario, deliberately scored low so the raw and final verdict are
    both WEAK, RECOMMEND COMPETITIVE ALTERNATIVE. Exercises: the
    ranked-alternatives DOCX/PPTX branch, the empty-mitigations-allowed
    path, and the NEW-mode recommended_next_action template."""
    return {
        "request": {
            "supplier": "BrightLeaf Creative Studio (illustrative)",
            "need_description": "Ongoing brand video production and motion graphics for the North America marketing team (illustrative demo scenario)",
            "stated_rationale": "Marketing says BrightLeaf is the only shop that gets the brand voice; no one else was seriously considered.",
            "mode": "NEW",
            "date": "July 22, 2026",
            "category": "Marketing Services, Video Production (illustrative)",
            "requester": "NA Marketing Communications (illustrative)",
            "est_value_usd": 240000,
            "term": "1 year, renewable annually",
        },
        "dimensions": [
            {"dimension_key": "unique_capability", "score": 1.5, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "Marketing team's own account of past work; no comparative review on file",
             "rationale": "The team describes BrightLeaf's style as distinctive, but distinctive style is not a structural capability; several agencies plausibly offer comparable creative work."},
            {"dimension_key": "constraint_basis", "score": 1.0, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "No structural constraint documentation provided",
             "rationale": "No IP, safety, regulatory, or continuity dependency identified. This reads as a stylistic preference, not a structural lock-in."},
            {"dimension_key": "competition_history", "score": 0.5, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "No RFP, RFI, or informal quote comparison on file (M365 search, this run, 0 results)",
             "rationale": "No competitive process of any kind has ever been run for this scope."},
            {"dimension_key": "requirements_separability", "score": 2.5, "label": "ASSERTED", "confidence": "MEDIUM",
             "evidence": "Marketing team description only",
             "rationale": "Ongoing production work is plausibly separable into a competed roster rather than a single sole-source shop, but this was never evaluated."},
            {"dimension_key": "alt_availability", "score": 2.0, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "Light market-check research log, this run (2 searches)",
             "rationale": "The light market-check surfaced at least 3 comparably-sized creative studios serving similar brand accounts; none were seriously evaluated."},
            {"dimension_key": "urgency_legitimacy", "score": 1.5, "label": "VERIFIED", "confidence": "MEDIUM",
             "evidence": "Marketing team's own project calendar (uploaded)",
             "rationale": "The urgency is self-created: the current statement of work was allowed to lapse and was renewed without a lead-time buffer."},
            {"dimension_key": "price_validation", "score": 1.0, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "No should-cost model or market-rate benchmark on file",
             "rationale": "The price is taken on faith; no should-cost model, market-rate benchmark, or rate card comparison exists for this spend."},
        ],
        "mitigations": [],
        "alternatives": [
            {"candidate_name": "Palmwood Studio Collective (illustrative)", "origin": "market-check",
             "capability_gap": "Smaller team, no dedicated NA account lead yet, but a comparable reel and a 20% lower day rate.",
             "reassessed_viability": "Viable", "confidence": "MEDIUM", "source": "Web market-check, this run", "date": "Jul 2026"},
            {"candidate_name": "Anchor and Ash Motion (illustrative)", "origin": "market-check",
             "capability_gap": "Strong motion graphics portfolio, limited live-action production experience.",
             "reassessed_viability": "Viable With Gaps", "confidence": "LOW", "source": "Web market-check, this run", "date": "Jul 2026"},
            {"candidate_name": "Tidewater Films (illustrative)", "origin": "user-provided",
             "capability_gap": "Marketing team's own shortlist candidate from a prior pitch; not evaluated further this run.",
             "reassessed_viability": "Viable With Gaps", "confidence": "LOW", "source": "Marketing team, this run", "date": "Jul 2026"},
        ],
        "price_validation": {"method": "none", "source": "", "band_low": None, "band_high": None, "sole_source_price": None, "position": "not_computable"},
        "research_log": [
            {"query": "brand video production agency alternative to BrightLeaf Creative Studio", "source": "Web search", "date": "Jul 20, 2026", "results": 8},
            {"query": "marketing video production vendor comparison North America 2026", "source": "Web search", "date": "Jul 20, 2026", "results": 6},
        ],
        "consumed_artifacts": [],
        "sme_routing": [],
        "research_methodology_note": (
            "No supplier-landscape or should-cost-builder output existed for this spend; the light "
            "market-check (2 searches) is the only alternatives evidence and is disclosed as RESEARCH "
            "PENDING for a full shortlist."
        ),
    }


def _demo_defensible_capped_register() -> Dict[str, Any]:
    """ILLUSTRATIVE DEMO DATA, a fictional AUDIT-mode calibration services
    scenario. Scored to land a RAW Defensible score (>= 4.0) but with a
    majority of dimensions labeled ASSERTED (carried forward from a prior
    justification memo, not independently re-verified this audit cycle),
    to exercise the Hard Rule 2 cap: the FINAL verdict must be forced down
    to DEFENSIBLE WITH MITIGATIONS with an auto-appended 'verify the
    ASSERTED claims' mitigation. Also exercises the RESEARCH PENDING path
    (only 1 research-log entry) and the gating-keyword-triggers-SME-routing
    invariant (the need description references GxP-adjacent equipment)."""
    return {
        "request": {
            "supplier": "Solstice Metrology Instruments (illustrative)",
            "need_description": "Calibration services for critical GxP measurement equipment across 2 sites (illustrative demo scenario)",
            "stated_rationale": "The original sole-source justification says Solstice's calibration standards are uniquely traceable to the site's validated baseline; being re-tested here as part of a Compliance-initiated audit.",
            "mode": "AUDIT",
            "date": "July 22, 2026",
            "category": "Lab and Clinical, Calibration Services (illustrative)",
            "requester": "Corporate Compliance, Audit Team (illustrative)",
            "est_value_usd": 610000,
            "term": "2 years, calibration services agreement",
        },
        "dimensions": [
            {"dimension_key": "unique_capability", "score": 4.5, "label": "ASSERTED", "confidence": "MEDIUM",
             "evidence": "Prior audit narrative, not independently re-verified this cycle",
             "rationale": "The prior justification asserts Solstice's calibration standards are uniquely traceable to the original validated baseline; this claim has not been independently re-verified in this audit cycle."},
            {"dimension_key": "constraint_basis", "score": 4.5, "label": "ASSERTED", "confidence": "MEDIUM",
             "evidence": "Prior justification memo, not independently re-verified this cycle",
             "rationale": "A continuity dependency on the original calibration chain is asserted in the prior justification memo but not independently re-confirmed this audit."},
            {"dimension_key": "competition_history", "score": 4.0, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "Formal RFP on file from original award, 2021 (uploaded)",
             "rationale": "A formal RFP was run at original award in 2021 and Solstice won on a documented technical and commercial basis."},
            {"dimension_key": "requirements_separability", "score": 3.5, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "Prior justification memo only, not reviewed this audit",
             "rationale": "The prior memo asserts the calibration chain cannot be split; this audit did not independently re-examine whether any portion is separable."},
            {"dimension_key": "alt_availability", "score": 4.0, "label": "ASSERTED", "confidence": "LOW",
             "evidence": "Prior justification memo's market note, not independently re-verified this cycle",
             "rationale": "The prior memo states no viable alternative exists; this audit did not run its own market-check to confirm that claim still holds four years later."},
            {"dimension_key": "urgency_legitimacy", "score": 4.5, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "Existing calibration services agreement end date (uploaded)",
             "rationale": "The agreement's expiration date is fixed and external; this is a standard renewal-timeline driven review, not a self-created urgency."},
            {"dimension_key": "price_validation", "score": 4.5, "label": "VERIFIED", "confidence": "HIGH",
             "evidence": "market_rate_benchmark.csv (market-rate-benchmarking output, consumed this run)",
             "rationale": "A market-rate benchmark for GxP calibration services independently validates the current rate as within the observed band."},
        ],
        "mitigations": [
            {"dimension": "Requirements Separability",
             "action": "Independently re-examine whether any portion of the calibration chain can be separated and competed, rather than relying on the 2021 memo's assertion.",
             "owner": "Category lead, Lab and Clinical", "due": "Before next renewal"},
        ],
        "alternatives": [
            {"candidate_name": "No alternatives evaluated this audit cycle", "origin": "user-provided",
             "original_exclusion_reason": "n/a, no re-scan performed",
             "capability_gap": "n/a, placeholder row: the original 2021 RFP evaluated 3 vendors but none were re-examined in this audit; a fresh market-check was out of scope for this compliance audit.",
             "reassessed_viability": "Not Viable", "confidence": "LOW", "source": "Audit scope note, this run", "date": "Jul 22, 2026"},
        ],
        "price_validation": {
            "method": "market-rate", "source": "market_rate_benchmark.csv (market-rate-benchmarking, consumed this run)",
            "band_low": 560000, "band_high": 650000, "sole_source_price": 610000, "position": "within",
        },
        "research_log": [
            {"query": "Review of prior 2021 RFP file and justification memo", "source": "Internal file review", "date": "Jul 22, 2026", "results": 1},
        ],
        "consumed_artifacts": [
            "market_rate_benchmark.csv (market-rate-benchmarking, price validation)",
            "prior sole-source justification memo, 2021 (Project Knowledge)",
        ],
        "sme_routing": [
            {"issue": "GxP-adjacent calibration equipment qualification status under the renewed agreement",
             "route_to": "Quality / GxP",
             "reason": "Any change to the calibration chain touches validated equipment; Compliance/Quality must clear qualification impact, this skill only flags it."},
        ],
        "research_methodology_note": (
            "This audit reviewed the prior 2021 RFP file and justification memo only; no fresh web "
            "market-check was run this cycle, which is disclosed as RESEARCH PENDING below."
        ),
    }


# ===========================================================================
# Self-test / CLI
# ===========================================================================

def _pptx_all_text(prs) -> str:
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
    return "\n".join(texts)


def _docx_all_text(doc) -> str:
    texts = [p.text for p in doc.paragraphs]
    for tbl in doc.tables:
        for row in tbl.rows:
            texts.extend(c.text for c in row.cells)
    return "\n".join(texts)


def _run_self_test() -> int:
    import tempfile
    import zipfile

    print("=" * 78)
    print("sole_source_generator.py self-test")
    print("=" * 78)
    print(f"numeric_kernel.py available: {KERNEL_AVAILABLE}")
    print(f"python-docx available: {DOCX_AVAILABLE}" + (f" (version {DOCX_VERSION})" if DOCX_AVAILABLE else ""))
    print(f"python-pptx available: {PPTX_AVAILABLE}" + (f" (version {PPTX_VERSION})" if PPTX_AVAILABLE else ""))
    print()

    results: List[tuple] = []

    def check(label, condition, detail=""):
        results.append((label, bool(condition), detail))
        status = "PASS" if condition else "FAIL"
        line = f"[{status}] {label}"
        if detail:
            line += f"  ({detail})"
        print(line)

    demo_a = _demo_defensible_with_mitigations_register()
    demo_b = _demo_weak_register()
    demo_c = _demo_defensible_capped_register()

    # --- Step 0: the three demo registers validate and compute correctly --
    try:
        reg_a = validate_sole_source_input(copy.deepcopy(demo_a))
        check("validate_sole_source_input accepts demo A (defensible-with-mitigations)", True)
    except Exception as e:
        check("validate_sole_source_input accepts demo A (defensible-with-mitigations)", False, str(e))
        raise

    gt_a = compute_ground_truth(reg_a)
    check(
        "Demo A Defensibility Score matches SKILL.md Phase 4's own worked example (3.475)",
        abs(gt_a.defensibility_score - 3.475) < 1e-9,
        f"got {gt_a.defensibility_score}",
    )
    check("Demo A verdict is DEFENSIBLE_WITH_MITIGATIONS", gt_a.verdict == "DEFENSIBLE_WITH_MITIGATIONS")
    check("Demo A verdict is NOT capped (majority VERIFIED, only 1 of 7 ASSERTED)", gt_a.verdict_capped is False)
    check("Demo A label counts: 6 VERIFIED, 1 ASSERTED, 0 INFERRED",
          gt_a.label_counts == {"VERIFIED": 6, "ASSERTED": 1, "INFERRED": 0}, f"got {gt_a.label_counts}")
    check("Demo A research_pending is False (3 research-log entries >= G7 minimum 2)", gt_a.research_pending is False)
    try:
        run_hardcoded_invariant_checks(reg_a, gt_a)
        check("All hard-coded invariant checks PASS on demo A", True)
    except ReconciliationError as e:
        check("All hard-coded invariant checks PASS on demo A", False, str(e))
        raise

    reg_b = validate_sole_source_input(copy.deepcopy(demo_b))
    gt_b = compute_ground_truth(reg_b)
    check(
        "Demo B Defensibility Score matches hand-derived expectation (1.4)",
        abs(gt_b.defensibility_score - 1.4) < 1e-9,
        f"got {gt_b.defensibility_score}",
    )
    check("Demo B verdict is WEAK_RECOMMEND_COMPETITION", gt_b.verdict == "WEAK_RECOMMEND_COMPETITION")
    check(
        "Demo B is majority-ASSERTED but NOT capped (cap only demotes a raw-Defensible score)",
        gt_b.asserted_majority is True and gt_b.verdict_capped is False,
        f"asserted_majority={gt_b.asserted_majority}, verdict_capped={gt_b.verdict_capped}",
    )
    check("Demo B has 0 effective mitigations (WEAK verdict does not require any)", len(gt_b.effective_mitigations) == 0)
    check(
        "Demo B recommended_next_action uses the NEW-mode WEAK-verdict template",
        "re-scope to separate the truly-constrained portion" in gt_b.recommended_next_action,
        gt_b.recommended_next_action[:120],
    )
    try:
        run_hardcoded_invariant_checks(reg_b, gt_b)
        check("All hard-coded invariant checks PASS on demo B", True)
    except ReconciliationError as e:
        check("All hard-coded invariant checks PASS on demo B", False, str(e))
        raise

    reg_c = validate_sole_source_input(copy.deepcopy(demo_c))
    gt_c = compute_ground_truth(reg_c)
    check(
        "Demo C raw Defensibility Score is >= 4.0 (raw DEFENSIBLE)",
        gt_c.raw_verdict == "DEFENSIBLE",
        f"score={gt_c.defensibility_score}, raw_verdict={gt_c.raw_verdict}",
    )
    check(
        "Demo C is majority ASSERTED (4 of 7) and IS capped to DEFENSIBLE_WITH_MITIGATIONS",
        gt_c.asserted_majority is True and gt_c.verdict == "DEFENSIBLE_WITH_MITIGATIONS" and gt_c.verdict_capped is True,
        f"label_counts={gt_c.label_counts}, verdict={gt_c.verdict}, capped={gt_c.verdict_capped}",
    )
    check(
        "Demo C auto-appends the 'verify the ASSERTED claims' mitigation (Hard Rule 2)",
        len(gt_c.effective_mitigations) == len(reg_c.mitigations) + 1
        and gt_c.effective_mitigations[-1].auto_generated is True,
        f"n_effective={len(gt_c.effective_mitigations)}, n_input={len(reg_c.mitigations)}",
    )
    check("Demo C research_pending is True (only 1 research-log entry < G7 minimum 2)", gt_c.research_pending is True)
    try:
        run_hardcoded_invariant_checks(reg_c, gt_c)
        check("All hard-coded invariant checks PASS on demo C", True)
    except ReconciliationError as e:
        check("All hard-coded invariant checks PASS on demo C", False, str(e))
        raise

    # --- Step 1: validation refusal tests ----------------------------------
    def _expect_validation_error(label, broken_register):
        try:
            validate_sole_source_input(broken_register)
            check(label, False, "did not raise")
        except SoleSourceValidationError as e:
            check(label, True, str(e)[:160])

    b1 = copy.deepcopy(demo_a)
    del b1["request"]["stated_rationale"]
    _expect_validation_error("refuses a register missing the blocking 'request.stated_rationale' field", b1)

    b2 = copy.deepcopy(demo_a)
    b2["dimensions"][0]["score"] = 7.5
    _expect_validation_error("refuses a dimension score outside [0.0, 5.0]", b2)

    b3 = copy.deepcopy(demo_a)
    del b3["dimensions"][0]
    _expect_validation_error("refuses a register with fewer than 7 dimensions", b3)

    b4 = copy.deepcopy(demo_a)
    b4["dimensions"][0]["dimension_key"] = "not_a_real_dimension"
    _expect_validation_error("refuses an unknown dimension_key", b4)

    b5 = copy.deepcopy(demo_a)
    b5["dimensions"][0]["weight"] = 0.99
    _expect_validation_error("refuses a caller-supplied weight that disagrees with the fixed weight table", b5)

    b6 = copy.deepcopy(demo_a)
    b6["dimensions"][0]["label"] = "PROBABLY"
    _expect_validation_error("refuses an unknown evidence label enum", b6)

    b7 = copy.deepcopy(demo_a)
    b7["alternatives"] = []
    _expect_validation_error("refuses an empty 'alternatives' list", b7)

    b8 = copy.deepcopy(demo_a)
    b8["price_validation"]["method"] = "none"
    _expect_validation_error("refuses method 'none' with populated band/price figures still present", b8)

    b9 = copy.deepcopy(demo_a)
    b9["price_validation"]["band_low"] = 500000  # now band_low > band_high (445000)
    _expect_validation_error("refuses band_low > band_high", b9)

    b10 = copy.deepcopy(demo_a)
    b10["request"]["need_description"] += " " + chr(0x2014) + " flagged for the em dash test."
    _expect_validation_error("refuses an em dash anywhere in the register", b10)

    b11 = copy.deepcopy(demo_a)
    b11["alternatives"][2]["origin"] = "supplier-landscape-excluded"
    b11["alternatives"][2]["original_exclusion_reason"] = ""
    _expect_validation_error(
        "refuses origin 'supplier-landscape-excluded' with an empty original_exclusion_reason", b11
    )

    # --- Step 2: reconciliation-invariant injection tests -------------------
    def _expect_reconciliation_error(label, fn):
        try:
            fn()
            check(label, False, "did not raise, but should have")
        except ReconciliationError as e:
            check(label, True, str(e)[:160])

    bad_gt1 = copy.deepcopy(gt_a)
    bad_gt1.weighted_contribution["unique_capability"] += 1.0
    _expect_reconciliation_error(
        "Weighted-contributions-reconcile invariant CORRECTLY REJECTS a corrupted contribution",
        lambda: _assert_weighted_contributions_reconcile(bad_gt1),
    )

    bad_gt2 = copy.deepcopy(gt_a)
    bad_gt2.defensibility_score = 9.0
    _expect_reconciliation_error(
        "Defensibility-score-range invariant CORRECTLY REJECTS an out-of-range score",
        lambda: _assert_defensibility_score_range(bad_gt2),
    )

    bad_gt3 = copy.deepcopy(gt_a)
    bad_gt3.raw_verdict = "WEAK_RECOMMEND_COMPETITION"  # disagrees with its own score (3.475)
    _expect_reconciliation_error(
        "Verdict-matches-score-band invariant CORRECTLY REJECTS a raw_verdict that disagrees with the score",
        lambda: _assert_verdict_matches_score_band_and_cap(bad_gt3),
    )

    bad_gt4 = copy.deepcopy(gt_c)
    bad_gt4.verdict = "DEFENSIBLE"  # silently un-cap a majority-ASSERTED case
    bad_gt4.verdict_capped = False
    _expect_reconciliation_error(
        "Rule-2-cap invariant CORRECTLY REJECTS silently un-capping a majority-ASSERTED Defensible score",
        lambda: _assert_verdict_matches_score_band_and_cap(bad_gt4),
    )

    bad_reg1 = copy.deepcopy(reg_a)
    bad_reg1.price_validation.position = "above"  # actually within band
    _expect_reconciliation_error(
        "Price-position-consistent invariant CORRECTLY REJECTS a position that disagrees with the figures",
        lambda: _assert_price_position_consistent(bad_reg1),
    )

    bad_gt5 = copy.deepcopy(gt_a)
    bad_gt5.verdict = "DEFENSIBLE_WITH_MITIGATIONS"
    bad_gt5.effective_mitigations = []
    _expect_reconciliation_error(
        "Mitigations-required-when-conditioned invariant CORRECTLY REJECTS a conditioned verdict with no mitigations",
        lambda: _assert_mitigations_required_when_conditioned(bad_gt5),
    )

    bad_reg2 = copy.deepcopy(reg_a)
    bad_reg2.alternatives = []
    _expect_reconciliation_error(
        "Alternatives-register-nonempty invariant CORRECTLY REJECTS an empty alternatives list",
        lambda: _assert_alternatives_nonempty(bad_reg2),
    )

    bad_gt6 = copy.deepcopy(gt_a)
    bad_gt6.ranked_alternatives = list(reversed(bad_gt6.ranked_alternatives))
    _expect_reconciliation_error(
        "Ranked-alternatives-sorted invariant CORRECTLY REJECTS a reversed ranking",
        lambda: _assert_ranked_alternatives_sorted(bad_gt6),
    )

    bad_reg3 = copy.deepcopy(reg_c)
    bad_reg3.sme_routing = []  # need_description references GxP; routing removed
    _expect_reconciliation_error(
        "Gating-items-routed invariant CORRECTLY REJECTS a GxP-referencing case with empty sme_routing",
        lambda: _assert_gating_items_routed(bad_reg3),
    )

    # --- Step 3: build real files (both formats, all three scenarios) -------
    if not DOCX_AVAILABLE or not PPTX_AVAILABLE:
        check(
            "python-docx and python-pptx both available to write real files", False,
            f"DOCX_AVAILABLE={DOCX_AVAILABLE}, PPTX_AVAILABLE={PPTX_AVAILABLE}",
        )
        print()
        print("Cannot proceed past this point without both python-docx and python-pptx.")
    else:
        check("python-docx and python-pptx both available to write real files", True,
              f"docx {DOCX_VERSION}, pptx {PPTX_VERSION}")

        tmp_dir = tempfile.gettempdir()
        scenarios = [
            ("A_defensible_with_mitigations", demo_a),
            ("B_weak", demo_b),
            ("C_capped", demo_c),
        ]
        built: Dict[str, Dict[str, str]] = {}
        for name, data in scenarios:
            built[name] = {}
            for fmt in ("docx", "pptx"):
                path = os.path.join(tmp_dir, f"sole_source_challenge_selftest_{name}.{fmt}")
                try:
                    generate_sole_source_challenge(data, path, fmt)
                    check(f"generate_sole_source_challenge() ran end-to-end for {name} ({fmt})", True)
                except Exception as e:
                    check(f"generate_sole_source_challenge() ran end-to-end for {name} ({fmt})", False, str(e))
                    raise
                built[name][fmt] = path
                exists = os.path.exists(path)
                size = os.path.getsize(path) if exists else 0
                check(f"{name} {fmt.upper()} file written, size={size} bytes", exists and size > 0)
                check(f"{name} {fmt.upper()} unzips cleanly (valid OOXML zip container)", zipfile.is_zipfile(path))

        # --- Re-open and structurally spot-check demo A (both formats) -----
        expected_docx_headings = [
            "01 Request Summary", "02 The Challenge", "03 Evidence and Price Validation",
            "04 Verdict and Recommendation", "05 Next Steps and SME Routing",
            "06 Research Methodology", "Appendix: Raw Scorecard",
        ]
        try:
            doc_a = Document(built["A_defensible_with_mitigations"]["docx"])
            text_a = _docx_all_text(doc_a)
            check("Re-opened demo A DOCX contains all 7 section headings (01-06 + Appendix)",
                  all(h in text_a for h in expected_docx_headings))
            check("Re-opened demo A DOCX names the supplier (Vaisala viewLinc)", "Vaisala viewLinc" in text_a)
            check("Re-opened demo A DOCX shows the Defensibility Score (3.475, 3-decimal appendix figure)", "3.475" in text_a)
            check("Re-opened demo A DOCX shows the Defensible With Mitigations verdict label",
                  "Defensible With Mitigations" in text_a)
            check("Re-opened demo A DOCX includes the Mitigations table heading", "Mitigations (condition the justification" in text_a)
            check("Re-opened demo A DOCX includes the Alternatives Considered subsection", "Alternatives considered" in text_a)
            check("Re-opened demo A DOCX includes the SME routing table", "Quality / GxP" in text_a)
            n_tables_a_docx = len(doc_a.tables)
            check("Re-opened demo A DOCX contains multiple tables", n_tables_a_docx >= 6, f"n_tables={n_tables_a_docx}")
        except Exception as e:
            check("Re-opened demo A DOCX structural spot-checks", False, str(e))

        try:
            prs_a = Presentation(built["A_defensible_with_mitigations"]["pptx"])
            text_a_pptx = _pptx_all_text(prs_a)
            check("Re-opened demo A PPTX has exactly 13 slides (1 title + 01..06 sections + Appendix)",
                  len(prs_a.slides) == 13, f"n_slides={len(prs_a.slides)}")
            check("Re-opened demo A PPTX names the supplier (Vaisala viewLinc) on the title slide",
                  "Vaisala viewLinc" in text_a_pptx)
            check("Re-opened demo A PPTX shows the Defensibility Score (3.47 or 3.475)",
                  "3.47" in text_a_pptx or "3.475" in text_a_pptx)
            check("Re-opened demo A PPTX shows the Defensible With Mitigations verdict label",
                  "Defensible With Mitigations" in text_a_pptx)
            check("Re-opened demo A PPTX includes the Mitigations slide", "Mitigations" in text_a_pptx)
            check("Re-opened demo A PPTX includes the Appendix, Raw Scorecard slide", "Raw Scorecard" in text_a_pptx)
            n_tables_a_pptx = sum(1 for s in prs_a.slides for sh in s.shapes if sh.has_table)
            check("Re-opened demo A PPTX contains multiple tables", n_tables_a_pptx >= 6, f"n_tables={n_tables_a_pptx}")
        except Exception as e:
            check("Re-opened demo A PPTX structural spot-checks", False, str(e))

        # --- Re-open and spot-check demo B (WEAK branch) --------------------
        try:
            doc_b = Document(built["B_weak"]["docx"])
            text_b = _docx_all_text(doc_b)
            check("Re-opened demo B DOCX shows the WEAK verdict label", "Weak, Recommend Competitive Alternative" in text_b)
            check("Re-opened demo B DOCX includes the ranked-alternatives heading (WEAK branch)",
                  "Ranked alternatives worth pursuing" in text_b)
            check("Re-opened demo B DOCX names the top-ranked alternative (Palmwood Studio Collective)",
                  "Palmwood Studio Collective" in text_b)
        except Exception as e:
            check("Re-opened demo B DOCX structural spot-checks", False, str(e))

        try:
            prs_b = Presentation(built["B_weak"]["pptx"])
            text_b_pptx = _pptx_all_text(prs_b)
            check("Re-opened demo B PPTX includes the Ranked Alternatives Worth Pursuing slide",
                  "Ranked Alternatives Worth Pursuing" in text_b_pptx)
            check("Re-opened demo B PPTX shows price validation method 'none' handled without a crash",
                  "not_computable" in text_b_pptx)
        except Exception as e:
            check("Re-opened demo B PPTX structural spot-checks", False, str(e))

        # --- Re-open and spot-check demo C (capped verdict) -----------------
        try:
            doc_c = Document(built["C_capped"]["docx"])
            text_c = _docx_all_text(doc_c)
            check("Re-opened demo C DOCX shows the capped Defensible With Mitigations verdict",
                  "Defensible With Mitigations" in text_c and "Capped from a raw Defensible score" in text_c)
            check("Re-opened demo C DOCX includes the auto-generated 'verify the ASSERTED claims' mitigation",
                  "Verify the ASSERTED claims" in text_c)
            check("Re-opened demo C DOCX includes the GxP SME routing row", "Quality / GxP" in text_c)
            check("Re-opened demo C DOCX flags RESEARCH PENDING (only 1 research-log entry)", "RESEARCH PENDING" in text_c)
        except Exception as e:
            check("Re-opened demo C DOCX structural spot-checks", False, str(e))

        try:
            prs_c = Presentation(built["C_capped"]["pptx"])
            text_c_pptx = _pptx_all_text(prs_c)
            check("Re-opened demo C PPTX shows the capped verdict narrative on the Verdict slide",
                  "Capped from a raw Defensible score" in text_c_pptx)
            check("Re-opened demo C PPTX includes the auto-generated mitigation on the Mitigations slide",
                  "Verify the ASSERTED claims" in text_c_pptx)
        except Exception as e:
            check("Re-opened demo C PPTX structural spot-checks", False, str(e))

    print()
    print("=" * 78)
    passed = sum(1 for _, ok, _ in results if ok)
    failed = sum(1 for _, ok, _ in results if not ok)
    total = len(results)
    print(f"SUMMARY: {passed}/{total} passed, {failed}/{total} failed")
    if failed:
        print("FAILED CASES:")
        for label, ok, detail in results:
            if not ok:
                print(f"  - {label}: {detail}")
    print("=" * 78)
    return 1 if failed else 0


def main(argv: Optional[List[str]] = None) -> int:
    """CLI entry point. Usage:
        python sole_source_generator.py --input register.json --output sole_source_challenge_report.docx --format docx
        python sole_source_generator.py --input register.json --output sole_source_challenge_report.pptx --format pptx
        python sole_source_generator.py --demo
        python sole_source_generator.py --self-test
        python sole_source_generator.py                 (no args -> runs the self-test)
    """
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate the sole_source_challenge_report (PPTX or DOCX) from a validated sole-source "
                    "register (sole-source-challenge-1c344a). See the module docstring for the input JSON schema."
    )
    parser.add_argument("--input", "-i", help="Path to a JSON file containing the sole-source register.")
    parser.add_argument("--output", "-o", help="Output file path (.docx or .pptx, matching --format).")
    parser.add_argument("--format", "-f", choices=["pptx", "docx"], default=None,
                         help="Output format: 'pptx' for the deck or 'docx' for the memo. Required with --input.")
    parser.add_argument("--demo", action="store_true",
                         help="Run the built-in self-test suite (generates all 3 illustrative demo registers "
                              "in both formats, reopens each, and asserts every expected section/slide/value).")
    parser.add_argument("--self-test", action="store_true", dest="self_test",
                         help="Alias for --demo.")
    args = parser.parse_args(argv)

    if args.demo or args.self_test or not args.input:
        return _run_self_test()

    if not args.output:
        print("ERROR: --output is required with --input.", file=sys.stderr)
        return 1
    if not args.format:
        print("ERROR: --format {pptx,docx} is required with --input.", file=sys.stderr)
        return 1

    import json
    try:
        with open(args.input, "r", encoding="utf-8") as f:
            raw = json.load(f)
    except (OSError, ValueError) as e:
        print(f"ERROR: could not read/parse --input {args.input!r}: {e}", file=sys.stderr)
        return 1

    try:
        generate_sole_source_challenge(raw, args.output, args.format)
    except ImportError as e:
        print(f"ERROR: {e}", file=sys.stderr)
        print("The sole-source register was valid and ground truth computed cleanly; only the "
              f"{args.format.upper()}-writing step could not run because the library is missing.", file=sys.stderr)
        return 1
    except (SoleSourceValidationError, ReconciliationError) as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 1

    print(f"Wrote {args.output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
